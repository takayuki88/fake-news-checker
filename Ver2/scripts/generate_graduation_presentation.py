from __future__ import annotations

import json
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DATASET_PATH = ROOT / "testdata" / "real_article_dataset.json"
EVAL_PATH = ROOT / "eval_real.json"
PLOTS_DIR = ROOT / "eval_real_plots"
OUTPUT_PATH = ROOT.parent / "卒業発表会_フェイクニュースチェッカー_20260329.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(247, 244, 238)
INK = RGBColor(32, 43, 58)
SUBTLE = RGBColor(95, 105, 120)
ACCENT = RGBColor(16, 110, 112)
ACCENT_2 = RGBColor(188, 98, 44)
ACCENT_3 = RGBColor(62, 132, 87)
LINE = RGBColor(210, 206, 198)
WHITE = RGBColor(255, 255, 255)
WARN = RGBColor(199, 134, 32)
ALERT = RGBColor(176, 57, 57)

FONT_BODY = "Yu Gothic"
FONT_TITLE = "Yu Gothic"


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def set_slide_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_header(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.8), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT_TITLE
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = INK
    p.alignment = PP_ALIGN.LEFT
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.62), Inches(0.95), Inches(8.5), Inches(0.35))
        sub_tf = sub_box.text_frame
        sub_p = sub_tf.paragraphs[0]
        sub_run = sub_p.add_run()
        sub_run.text = subtitle
        sub_run.font.name = FONT_BODY
        sub_run.font.size = Pt(11)
        sub_run.font.color.rgb = SUBTLE
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(1.18), Inches(1.55), Inches(0.07)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def add_footer(slide, note: str = "Fake News Checker 2 Ver2 | 2026-03-29 時点の成果物より作成") -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.22))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = note
    run.font.name = FONT_BODY
    run.font.size = Pt(9)
    run.font.color.rgb = SUBTLE
    p.alignment = PP_ALIGN.RIGHT


def add_textbox(slide, left, top, width, height, text: str = "", font_size: int = 20, color=INK, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if text:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = FONT_BODY
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
    return box


def add_bullets(slide, left, top, width, height, items: list[str], font_size: int = 20, color=INK) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.space_after = Pt(8)
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = FONT_BODY
            run.font.size = Pt(font_size)
            run.font.color.rgb = color


def add_panel(slide, left, top, width, height, fill_rgb=WHITE, line_rgb=LINE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = line_rgb
    shape.line.width = Pt(1.2)
    return shape


def add_metric_card(slide, left, top, width, height, label: str, value: str, accent: RGBColor) -> None:
    card = add_panel(slide, left, top, width, height, fill_rgb=WHITE)
    card.line.color.rgb = accent
    accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.08))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()
    label_box = add_textbox(slide, left + Inches(0.18), top + Inches(0.20), width - Inches(0.35), Inches(0.25), label, 12, SUBTLE, False)
    value_box = add_textbox(slide, left + Inches(0.18), top + Inches(0.42), width - Inches(0.35), Inches(0.48), value, 24, INK, True)
    label_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    value_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT


def add_table_like(slide, left, top, width, row_h, rows: list[tuple[str, str]], title: str, accent: RGBColor) -> None:
    total_h = row_h * (len(rows) + 1)
    add_panel(slide, left, top, width, total_h, fill_rgb=WHITE)
    head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, row_h)
    head.fill.solid()
    head.fill.fore_color.rgb = accent
    head.line.fill.background()
    add_textbox(slide, left + Inches(0.16), top + Inches(0.04), width - Inches(0.3), row_h - Inches(0.05), title, 16, WHITE, True)
    for idx, (label, value) in enumerate(rows, start=1):
        y = top + row_h * idx
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left + Inches(0.01), y, width - Inches(0.02), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = LINE
        line.line.fill.background()
        add_textbox(slide, left + Inches(0.16), y + Inches(0.06), width * 0.55, row_h - Inches(0.08), label, 16, INK, False)
        box = add_textbox(slide, left + width * 0.62, y + Inches(0.06), width * 0.30, row_h - Inches(0.08), value, 16, INK, True)
        box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_process_box(slide, left, top, width, height, title: str, desc: str, fill_rgb: RGBColor) -> None:
    shape = add_panel(slide, left, top, width, height, fill_rgb=fill_rgb, line_rgb=fill_rgb)
    shape.line.color.rgb = fill_rgb
    add_textbox(slide, left + Inches(0.14), top + Inches(0.12), width - Inches(0.28), Inches(0.28), title, 16, WHITE, True)
    desc_box = add_textbox(slide, left + Inches(0.14), top + Inches(0.45), width - Inches(0.28), height - Inches(0.56), desc, 13, WHITE, False)
    desc_box.text_frame.vertical_anchor = MSO_ANCHOR.TOP


def add_connector(slide, left, top, width, height, color: RGBColor) -> None:
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()


def build_presentation() -> Presentation:
    dataset = load_json(DATASET_PATH)
    evaluation = load_json(EVAL_PATH)
    domain_counts = Counter(case["expected"]["domain"] for case in dataset["cases"])
    mismatch_counts = Counter((item["truth"], item["pred"]) for item in evaluation["mismatches"])
    top_mismatches = mismatch_counts.most_common(3)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    hero.fill.solid()
    hero.fill.fore_color.rgb = RGBColor(22, 40, 57)
    hero.line.fill.background()
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(0.95), Inches(0.23), Inches(4.9))
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT_2
    band.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(1.05), Inches(1.1), Inches(8.8), Inches(2.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "フェイクニュースチェッカー\n作成まとめ"
    r.font.name = FONT_TITLE
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE
    p.space_after = Pt(0)
    sub = slide.shapes.add_textbox(Inches(1.08), Inches(3.35), Inches(7.8), Inches(1.0))
    sub_tf = sub.text_frame
    sub_tf.word_wrap = True
    p = sub_tf.paragraphs[0]
    r = p.add_run()
    r.text = "卒業発表会向けプレゼン資料\nFake News Checker 2 Ver2 / 2026年3月29日時点"
    r.font.name = FONT_BODY
    r.font.size = Pt(18)
    r.font.color.rgb = RGBColor(218, 225, 232)
    badge = add_panel(slide, Inches(9.55), Inches(1.25), Inches(2.8), Inches(4.7), fill_rgb=RGBColor(242, 238, 229), line_rgb=RGBColor(242, 238, 229))
    badge.line.fill.background()
    add_textbox(slide, Inches(9.9), Inches(1.7), Inches(2.2), Inches(0.4), "発表で扱う内容", 16, SUBTLE, True)
    add_bullets(
        slide,
        Inches(9.75),
        Inches(2.15),
        Inches(2.25),
        Inches(3.2),
        [
            "課題設定",
            "使用データ",
            "仮説とアプローチ",
            "結果と考察",
            "ビジネスインパクト",
            "工夫点と今後",
        ],
        font_size=16,
        color=INK,
    )

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "1. 課題と選んだ理由", "拡散前の一次確認を支援する仕組みを作る")
    add_panel(slide, Inches(0.65), Inches(1.55), Inches(5.95), Inches(4.95))
    add_panel(slide, Inches(6.75), Inches(1.55), Inches(5.95), Inches(4.95), fill_rgb=RGBColor(237, 244, 243))
    add_textbox(slide, Inches(0.9), Inches(1.8), Inches(5.2), Inches(0.35), "背景", 20, ACCENT, True)
    add_bullets(
        slide,
        Inches(0.88),
        Inches(2.2),
        Inches(5.2),
        Inches(3.8),
        [
            "SNSやニュースでは、真実と誇張、文脈不足、完全な誤情報が同じ画面に並ぶ",
            "特に医療・災害・金融の話題は、誤情報の拡散コストが大きい",
            "真偽だけでなく『どこを確認すべきか』まで示す支援ツールが必要",
        ],
        font_size=20,
    )
    add_textbox(slide, Inches(7.0), Inches(1.8), Inches(5.0), Inches(0.35), "このテーマを選んだ理由", 20, ACCENT, True)
    add_bullets(
        slide,
        Inches(6.98),
        Inches(2.2),
        Inches(5.2),
        Inches(3.6),
        [
            "判定結果だけではなく、確認リンクや説明を出して行動につなげたかった",
            "人手確認を前提にしつつ、危険な投稿を早めに見つける補助にしたかった",
            "最終目標は『自動断定』ではなく『拡散前の立ち止まりを促す』こと",
        ],
        font_size=20,
    )
    add_metric_card(slide, Inches(7.05), Inches(5.45), Inches(1.6), Inches(0.82), "入力", "URL / 本文", ACCENT)
    add_metric_card(slide, Inches(8.85), Inches(5.45), Inches(1.5), Inches(0.82), "出力", "5区分", ACCENT_2)
    add_metric_card(slide, Inches(10.55), Inches(5.45), Inches(1.75), Inches(0.82), "補助情報", "要注意度・根拠", ACCENT_3)
    add_footer(slide)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "2. 使用したデータについて", "評価用データセットは 2026-03-29 curated の balanced 100 件")
    add_textbox(slide, Inches(0.75), Inches(1.55), Inches(12), Inches(0.45), "評価対象は記事全文ではなく、analysis_text にまとめた『中心命題』です。", 18, SUBTLE, False)
    add_table_like(
        slide,
        Inches(0.75),
        Inches(2.0),
        Inches(3.65),
        Inches(0.55),
        [("正確", "20"), ("ほぼ正確", "20"), ("判断保留", "20"), ("不正確", "20"), ("誤り", "20")],
        "5区分の件数",
        ACCENT,
    )
    add_table_like(
        slide,
        Inches(4.65),
        Inches(2.0),
        Inches(3.65),
        Inches(0.55),
        [("医療", str(domain_counts["医療"])), ("一般", str(domain_counts["一般"])), ("災害", str(domain_counts["災害"])), ("金融", str(domain_counts["金融"]))],
        "ドメイン内訳",
        ACCENT_2,
    )
    add_panel(slide, Inches(8.55), Inches(2.0), Inches(4.05), Inches(3.3), fill_rgb=RGBColor(237, 244, 243))
    add_textbox(slide, Inches(8.85), Inches(2.25), Inches(3.4), Inches(0.35), "データの作り方", 18, ACCENT, True)
    add_bullets(
        slide,
        Inches(8.8),
        Inches(2.58),
        Inches(3.35),
        Inches(2.25),
        [
            "正確・ほぼ正確: 公的ページやファクトチェック記事の要約",
            "判断保留・不正確・誤り: 境界評価用に curated したケース",
            "100件を均等配分にして、判定の偏りを見えやすくした",
        ],
        font_size=16,
    )
    add_panel(slide, Inches(0.75), Inches(5.55), Inches(11.85), Inches(1.1), fill_rgb=RGBColor(252, 249, 242))
    add_bullets(
        slide,
        Inches(0.95),
        Inches(5.78),
        Inches(11.4),
        Inches(0.6),
        [
            "公開 verdict の 5区分評価を目的に作られており、support は各ラベル 20 件で固定",
            "README とテストでは、件数とラベル分布を自動確認できるようにしている",
        ],
        font_size=18,
    )
    add_footer(slide)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "3. 初期仮説・分析アプローチ", "ヒューリスティックで一次判定し、Gemini で外部根拠比較と書き振り評価を重ねる設計")
    add_panel(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(2.15), fill_rgb=RGBColor(252, 249, 242))
    add_bullets(
        slide,
        Inches(0.95),
        Inches(1.95),
        Inches(11.5),
        Inches(1.4),
        [
            "仮説1: 出典の有無、著者・日付、リンク数、表現の強さなどで『危険度の目安』は作れる",
            "仮説2: 境界ケースは外部根拠との照合を入れることで精度改善が期待できる",
            "仮説3: 煽り表現の強さは真偽とは切り離して別カードで見せた方が誤判定を減らせる",
        ],
        font_size=18,
    )
    y = Inches(4.2)
    box_w = Inches(2.1)
    add_process_box(slide, Inches(0.75), y, box_w, Inches(1.45), "入力", "URL または本文を受け取り、必要なら本文抽出", ACCENT)
    add_connector(slide, Inches(2.97), y + Inches(0.47), Inches(0.38), Inches(0.42), RGBColor(174, 196, 196))
    add_process_box(slide, Inches(3.3), y, box_w, Inches(1.45), "一次判定", "ヒューリスティックで リスク / 確信度 / ドメインを推定", ACCENT_2)
    add_connector(slide, Inches(5.52), y + Inches(0.47), Inches(0.38), Inches(0.42), RGBColor(219, 189, 169))
    add_process_box(slide, Inches(5.85), y, box_w, Inches(1.45), "Gemini 比較", "google_search と url_context で外部根拠比較", ACCENT_3)
    add_connector(slide, Inches(8.07), y + Inches(0.47), Inches(0.38), Inches(0.42), RGBColor(183, 209, 190))
    add_process_box(slide, Inches(8.4), y, box_w, Inches(1.45), "書き振り評価", "煽り / 断定 / 陰謀論テンプレを別軸で採点", WARN)
    add_connector(slide, Inches(10.62), y + Inches(0.47), Inches(0.38), Inches(0.42), RGBColor(221, 204, 164))
    add_process_box(slide, Inches(10.95), y, Inches(1.35), Inches(1.45), "出力", "5区分・要注意度・理由・確認リンク", ALERT)
    add_footer(slide)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "4. 改善の流れ", "Ver1 から Ver2 で『ラベル整理』『評価基盤』『説明性』を強化")
    add_panel(slide, Inches(0.8), Inches(1.75), Inches(5.3), Inches(4.95))
    add_panel(slide, Inches(7.1), Inches(1.75), Inches(5.3), Inches(4.95), fill_rgb=RGBColor(237, 244, 243))
    add_textbox(slide, Inches(1.08), Inches(2.05), Inches(2.1), Inches(0.35), "Ver1", 22, ACCENT_2, True)
    add_bullets(
        slide,
        Inches(1.0),
        Inches(2.45),
        Inches(4.6),
        Inches(3.6),
        [
            "6分類: 信頼性が高い / おおむね正確 / 判断保留 / 誤解を招く可能性が高い / フェイクの可能性が高い / 未確認",
            "Gemini 比較はあるが、評価指標やデータセットが整理されていなかった",
            "判定ラベルが多く、発表時に結果を説明しにくかった",
        ],
        font_size=18,
    )
    add_textbox(slide, Inches(7.4), Inches(2.05), Inches(2.1), Inches(0.35), "Ver2", 22, ACCENT, True)
    add_bullets(
        slide,
        Inches(7.3),
        Inches(2.45),
        Inches(4.65),
        Inches(3.6),
        [
            "公開用 verdict を 5分類に再設計し、要注意度バンドも 0〜100% に統一",
            "書き振り評価を真偽から分離し、高スコア時だけ『人による確認推奨』へ寄せる",
            "dataset_runner / evaluation / plot を追加し、定量評価できる形にした",
            "URL 取得では hard block + robots.txt を既定にして安全面も整理",
        ],
        font_size=18,
    )
    add_connector(slide, Inches(6.2), Inches(3.45), Inches(0.55), Inches(0.55), ACCENT)
    add_textbox(slide, Inches(5.55), Inches(4.2), Inches(1.8), Inches(0.6), "評価可能な\nプロトタイプへ", 17, INK, True)
    add_footer(slide)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "5. 分析結果", "2026年3月29日時点の eval_real.json を要約")
    add_metric_card(slide, Inches(0.8), Inches(1.55), Inches(1.9), Inches(0.95), "Accuracy", "44%", ACCENT)
    add_metric_card(slide, Inches(2.95), Inches(1.55), Inches(2.0), Inches(0.95), "Macro F1", "0.382", ACCENT_2)
    add_metric_card(slide, Inches(5.2), Inches(1.55), Inches(2.3), Inches(0.95), "実行件数", "100 / 100", ACCENT_3)
    add_metric_card(slide, Inches(7.75), Inches(1.55), Inches(2.15), Inches(0.95), "Error", "0件", WARN)
    add_metric_card(slide, Inches(10.15), Inches(1.55), Inches(2.15), Inches(0.95), "モデル", "heuristic", ALERT)
    plot_path = PLOTS_DIR / "evaluation_dashboard.png"
    slide.shapes.add_picture(str(plot_path), Inches(0.82), Inches(2.75), width=Inches(7.1))
    add_panel(slide, Inches(8.2), Inches(2.75), Inches(4.1), Inches(3.3), fill_rgb=RGBColor(252, 249, 242))
    add_textbox(slide, Inches(8.5), Inches(3.0), Inches(3.2), Inches(0.35), "読み取れること", 18, ACCENT, True)
    add_bullets(
        slide,
        Inches(8.45),
        Inches(3.32),
        Inches(3.35),
        Inches(2.4),
        [
            "『正確』は高精度で拾えており、F1 は 0.974",
            "『ほぼ正確』は一定精度だが、境界ケースでぶれやすい",
            "この評価は Gemini 未使用で、ローカル一次判定のみの結果",
        ],
        font_size=17,
    )
    note = add_textbox(slide, Inches(8.48), Inches(5.75), Inches(3.4), Inches(0.38), "※ predictions_real.json の meta では use_gemini=false", 12, SUBTLE, False)
    note.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    add_footer(slide)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "6. 考察", "現在のヒューリスティックは『中間寄り』に倒れやすく、危険ケースを鋭く切れていない")
    slide.shapes.add_picture(str(PLOTS_DIR / "confusion_matrix.png"), Inches(0.78), Inches(1.75), width=Inches(5.4))
    add_panel(slide, Inches(6.45), Inches(1.75), Inches(5.85), Inches(4.95))
    add_textbox(slide, Inches(6.8), Inches(2.05), Inches(4.9), Inches(0.35), "主な誤判定パターン", 20, ACCENT_2, True)
    mismatch_items = [
        f"{truth} → {pred}: {count}件"
        for (truth, pred), count in top_mismatches
    ]
    add_bullets(slide, Inches(6.78), Inches(2.45), Inches(4.9), Inches(1.2), mismatch_items, font_size=18)
    add_textbox(slide, Inches(6.8), Inches(3.95), Inches(4.8), Inches(0.35), "考察", 20, ACCENT, True)
    add_bullets(
        slide,
        Inches(6.78),
        Inches(4.3),
        Inches(4.9),
        Inches(1.95),
        [
            "公的ソースを含む『安全側』の見極めは得意",
            "一方で『判断保留』を楽観的に見てしまい、『誤り』も『不正確』止まりになる",
            "今後は外部反証の重み付けと、灰色ケース専用ルールの強化が必要",
        ],
        font_size=18,
    )
    add_footer(slide)

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "7. ビジネスインパクト", "現時点では『自動判定ツール』より『確認行動を促す支援ツール』として価値が高い")
    add_metric_card(slide, Inches(0.85), Inches(1.75), Inches(3.65), Inches(1.1), "想定利用シーン", "投稿前のセルフチェック", ACCENT)
    add_metric_card(slide, Inches(4.85), Inches(1.75), Inches(3.65), Inches(1.1), "組織導入の価値", "一次仕分けの補助", ACCENT_2)
    add_metric_card(slide, Inches(8.85), Inches(1.75), Inches(3.45), Inches(1.1), "教育面の価値", "メディアリテラシー向上", ACCENT_3)
    add_panel(slide, Inches(0.85), Inches(3.15), Inches(11.45), Inches(3.1), fill_rgb=RGBColor(252, 249, 242))
    add_bullets(
        slide,
        Inches(1.15),
        Inches(3.48),
        Inches(10.9),
        Inches(2.4),
        [
            "入力した文章に対して、5区分判定だけでなく確認先リンクを出すことで『自分で確かめる』行動を支援できる",
            "学校・自治体・メディア運用では、高リスク投稿の候補抽出や人手確認の優先順位づけに使える",
            "ただし 2026年3月29日時点の定量結果では誤情報の鋭い検出が弱いため、単独の自動判定用途にはまだ不十分",
            "価値の中心は『断定の代行』ではなく『確認コストの削減と注意喚起』にある",
        ],
        font_size=18,
    )
    add_footer(slide)

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "8. 工夫点と今後さらに実施したいこと", "評価基盤を作れたことが次の改善サイクルにつながる")
    add_panel(slide, Inches(0.8), Inches(1.7), Inches(5.75), Inches(4.95))
    add_panel(slide, Inches(6.8), Inches(1.7), Inches(5.75), Inches(4.95), fill_rgb=RGBColor(237, 244, 243))
    add_textbox(slide, Inches(1.08), Inches(2.0), Inches(2.5), Inches(0.35), "工夫点", 20, ACCENT, True)
    add_bullets(
        slide,
        Inches(1.0),
        Inches(2.35),
        Inches(5.0),
        Inches(3.65),
        [
            "真偽判定と書き振り評価を分離し、煽りの強さだけで真偽を決めない設計にした",
            "5区分100件の balanced dataset を整備して、改善前後を比較できるようにした",
            "評価 JSON と可視化 PNG を自動生成できるようにして、発表に使える形へ落とし込んだ",
            "Gemini が使えない場合でもローカル判定で止まらないフォールバックを入れた",
        ],
        font_size=18,
    )
    add_textbox(slide, Inches(7.08), Inches(2.0), Inches(3.0), Inches(0.35), "今後やりたいこと", 20, ACCENT_2, True)
    add_bullets(
        slide,
        Inches(7.0),
        Inches(2.35),
        Inches(5.0),
        Inches(3.7),
        [
            "Gemini 有効時でも同じ 100 件を評価し、ヒューリスティック単体との差分を確認する",
            "『判断保留』と『誤り』の境界を改善するルールや反証スコアを追加する",
            "最新ニュースや SNS 投稿の実データに近いセットへ拡張して、時事性への強さを見る",
            "利用者のフィードバックを蓄積し、誤判定ケースから継続学習できる形にする",
        ],
        font_size=18,
    )
    close_box = add_panel(slide, Inches(0.8), Inches(6.25), Inches(11.75), Inches(0.5), fill_rgb=RGBColor(252, 249, 242))
    close_box.line.color.rgb = LINE
    closing = add_textbox(
        slide,
        Inches(1.05),
        Inches(6.35),
        Inches(11.2),
        Inches(0.25),
        "まとめ: 現段階では『説明できる一次チェック支援ツール』として成立し、定量評価を土台に次の改善へ進める状態まで到達できた。",
        16,
        INK,
        True,
    )
    closing.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_footer(slide)

    return prs


def main() -> None:
    prs = build_presentation()
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(OUTPUT_PATH)


if __name__ == "__main__":
    main()
