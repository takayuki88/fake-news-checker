from __future__ import annotations

import json
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PRESENTATION_DIR = Path(__file__).resolve().parent
PROJECT_ROOT = PRESENTATION_DIR.parent

DATASET_PATH = PROJECT_ROOT / "testdata" / "shared" / "real_article_dataset_v2.json"
VER4_EVAL_PATH = PROJECT_ROOT / "Ver4" / "evaluation_outputs" / "20260427-0304" / "eval_real_article_dataset_v2_use_gpt_gemini.json"
VER4_PLOTS_DIR = PROJECT_ROOT / "Ver4" / "evaluation_outputs" / "20260427-0304" / "plots"
VER2_EVAL_PATH = PROJECT_ROOT / "Ver2" / "evaluation_outputs" / "20260411-0509" / "eval_real_article_dataset_v2_use_gemini.json"
VER3_EVAL_PATH = PROJECT_ROOT / "Ver3" / "evaluation_outputs" / "20260421-2001" / "eval_real_article_dataset_v2_use_gemini.json"

OUTPUT_PPTX = PRESENTATION_DIR / "fake_news_checker_ver4_final_presentation_20260429.pptx"
OUTPUT_NOTES = PRESENTATION_DIR / "fake_news_checker_ver4_final_presentation_20260429_notes.md"

FONT = "Yu Gothic"
BG = RGBColor(247, 248, 246)
INK = RGBColor(27, 35, 43)
SUBTLE = RGBColor(90, 101, 112)
WHITE = RGBColor(255, 255, 255)
NAVY = RGBColor(24, 48, 72)
TEAL = RGBColor(15, 124, 128)
GREEN = RGBColor(61, 143, 99)
ORANGE = RGBColor(207, 115, 47)
RED = RGBColor(182, 63, 63)
GOLD = RGBColor(191, 151, 48)
MINT = RGBColor(230, 243, 240)
BLUE_SOFT = RGBColor(232, 239, 248)
ORANGE_SOFT = RGBColor(250, 238, 227)
GREEN_SOFT = RGBColor(232, 244, 235)
RED_SOFT = RGBColor(249, 232, 232)
LINE = RGBColor(214, 220, 224)


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def pct(value: float) -> str:
    return f"{value * 100:.1f}%"


def short(text: str, limit: int = 50) -> str:
    return text if len(text) <= limit else text[: limit - 1] + "…"


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_text(slide, left, top, width, height, text="", size=20, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(3)
    tf.margin_right = Pt(3)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if text:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        p.alignment = align
    return box


def add_header(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, Inches(0.6), Inches(0.32), Inches(9.7), Inches(0.44), title, 24, NAVY, True)
    if subtitle:
        add_text(slide, Inches(0.62), Inches(0.78), Inches(10.6), Inches(0.28), subtitle, 10, SUBTLE)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(1.08), Inches(12.1), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def add_footer(slide, page: int) -> None:
    add_text(slide, Inches(0.65), Inches(7.05), Inches(7.2), Inches(0.18), "Fake News Checker Ver4 | 分析アプローチ", 8, SUBTLE)
    add_text(slide, Inches(12.25), Inches(7.05), Inches(0.55), Inches(0.18), str(page), 8, SUBTLE, False, PP_ALIGN.RIGHT)


def add_panel(slide, left, top, width, height, fill=WHITE, line=LINE, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_bullets(slide, left, top, width, height, items: list[str], size=17, color=INK, gap=0):
    box = add_text(slide, left, top, width, height)
    tf = box.text_frame
    tf.clear()
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.level = 0
        p.space_after = Pt(gap)
        p.line_spacing = 1.12
    return box


def add_metric(slide, left, top, width, height, label: str, value: str, accent=TEAL, note: str | None = None):
    add_panel(slide, left, top, width, height)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.08))
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()
    add_text(slide, left + Inches(0.12), top + Inches(0.2), width - Inches(0.24), Inches(0.22), label, 11, SUBTLE, True, PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.12), top + Inches(0.48), width - Inches(0.24), Inches(0.45), value, 25, accent, True, PP_ALIGN.CENTER)
    if note:
        add_text(slide, left + Inches(0.12), top + Inches(0.95), width - Inches(0.24), Inches(0.26), note, 9, SUBTLE, False, PP_ALIGN.CENTER)


def add_step(slide, left, top, width, height, title: str, body: str, accent=TEAL, number: str | None = None):
    add_panel(slide, left, top, width, height)
    if number:
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.12), top + Inches(0.14), Inches(0.34), Inches(0.34))
        badge.fill.solid()
        badge.fill.fore_color.rgb = accent
        badge.line.fill.background()
        add_text(slide, left + Inches(0.12), top + Inches(0.19), Inches(0.34), Inches(0.16), number, 9, WHITE, True, PP_ALIGN.CENTER)
        title_left = left + Inches(0.54)
        title_width = width - Inches(0.65)
    else:
        title_left = left + Inches(0.16)
        title_width = width - Inches(0.32)
    add_text(slide, title_left, top + Inches(0.14), title_width, Inches(0.25), title, 13, accent, True)
    add_text(slide, left + Inches(0.16), top + Inches(0.48), width - Inches(0.32), height - Inches(0.58), body, 10, INK)


def add_arrow(slide, left, top, width, height, color=SUBTLE):
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow


def add_picture_fit(slide, image_path: Path, left, top, width, height):
    add_panel(slide, left, top, width, height, WHITE)
    pic = slide.shapes.add_picture(str(image_path), left + Inches(0.08), top + Inches(0.08))
    max_w = width - Inches(0.16)
    max_h = height - Inches(0.16)
    scale = min(max_w / pic.width, max_h / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = left + int((width - pic.width) / 2)
    pic.top = top + int((height - pic.height) / 2)


def add_bar(slide, left, top, width, height, label: str, value: float, accent=TEAL, max_value=1.0):
    add_text(slide, left, top - Inches(0.02), Inches(1.05), Inches(0.2), label, 10, INK, True)
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left + Inches(1.2), top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(229, 234, 238)
    bg.line.fill.background()
    filled = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left + Inches(1.2), top, int(width * min(value / max_value, 1)), height)
    filled.fill.solid()
    filled.fill.fore_color.rgb = accent
    filled.line.fill.background()
    add_text(slide, left + Inches(1.2) + width + Inches(0.12), top - Inches(0.03), Inches(0.65), Inches(0.2), pct(value), 9, SUBTLE)


def build_notes(ver4: dict, dataset: dict, mismatch_patterns: list[tuple[tuple[str, str], int]]) -> str:
    cases = dataset["cases"]
    verdict_counts = Counter(case["expected_verdict"] for case in cases)
    top_patterns = [f"{truth} -> {pred}: {count}件" for (truth, pred), count in mismatch_patterns[:5]]
    return "\n".join(
        [
            "# Fake News Checker Ver4 Final Presentation Notes",
            "",
            "- generated_on: 2026-04-29",
            f"- dataset: {DATASET_PATH.relative_to(PROJECT_ROOT)}",
            f"- ver4_eval: {VER4_EVAL_PATH.relative_to(PROJECT_ROOT)}",
            "",
            "## Slide 1",
            "- 発表全体の表紙。Fake News Checkerを、Ver4の検証結果を中心に説明する。",
            "",
            "## Slide 2",
            "- 課題と選んだ理由。SNSやニュースの誤情報は拡散が速く、人手確認だけでは追いつきにくい。",
            "- 個人的な関心として、LLMを情報確認支援に使えるかを検証した。",
            "",
            "## Slide 3",
            "- 使用データ。100件、5ラベル各20件のbalanced datasetを使った。",
            "",
            "## Slide 4",
            "- 初期仮説は、ニュース本文だけでは真偽判定が安定しないというもの。",
            "- 数字、時期、人物名、根拠不足、一部だけ正しい主張が混ざるため、二値分類だけでは粗い。",
            "",
            "## Slide 5",
            "- 入力から最終判定まで、ローカル一次判定、GPT一次レビュー、Gemini外部根拠確認、補正ルールの順に処理する。",
            "",
            "## Slide 6",
            "- 5段階ラベルにより、完全な正誤だけでなく、軽微な差分や判断保留を表現する。",
            "",
            "## Slide 7",
            f"- 評価データは100件、5ラベルが各{min(verdict_counts.values())}件のバランス構成。",
            f"- Accuracy {pct(ver4['multiclass']['accuracy'])}、Macro F1 {pct(ver4['multiclass']['macro_avg']['f1'])}。",
            "",
            "## Slide 8",
            "- 改善の中心はモデル再学習ではなく、LLM出力をどう5段階ラベルへ写像するかのキャリブレーション。",
            "",
            "## Slide 9",
            "- 残ミスマッチは10件。",
            *[f"- {pattern}" for pattern in top_patterns],
            "",
            "## Slide 10",
            "- 保存済み評価ではVer4が最良。ただし同じ100件に合わせているため、未知データでの確認が必要。",
            "",
            "## Slide 11",
            "- ビジネスインパクト。大量投稿の一次スクリーニング、確認優先度づけ、説明付きレポート作成に使える。",
            "",
            "## Slide 12",
            "- 工夫点。評価bundle化、5段階分類、境界判定ルールの改善を行った。",
            "",
            "## Slide 13",
            "- 今後さらに実施したいこと。",
            "- 次は新規blind 100件を作り、各バージョンを同じ条件で比較する。",
            "",
            "## Slide 14",
            "- まとめは、既存LLM、外部根拠確認、ルールベース補正を組み合わせた確認支援ツールという位置づけ。",
        ]
    ) + "\n"


def build_presentation() -> None:
    dataset = load_json(DATASET_PATH)
    verdict_counts = Counter(case["expected_verdict"] for case in dataset["cases"])
    domain_counts = Counter(case["expected_domain"] for case in dataset["cases"])
    ver4 = load_json(VER4_EVAL_PATH)
    ver2 = load_json(VER2_EVAL_PATH)
    ver3 = load_json(VER3_EVAL_PATH)
    mismatch_patterns = Counter((item["truth"], item["pred"]) for item in ver4["mismatches"]).most_common()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    cover_band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    cover_band.fill.solid()
    cover_band.fill.fore_color.rgb = NAVY
    cover_band.line.fill.background()
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    add_text(slide, Inches(0.82), Inches(1.25), Inches(9.7), Inches(0.42), "Fake News Checker", 25, RGBColor(205, 229, 227), True)
    add_text(slide, Inches(0.82), Inches(1.82), Inches(10.8), Inches(1.1), "最終発表資料", 48, WHITE, True)
    add_text(slide, Inches(0.86), Inches(3.05), Inches(8.9), Inches(0.58), "課題・データ・分析アプローチ・結果・ビジネスインパクト", 19, RGBColor(232, 238, 244))
    add_metric(slide, Inches(0.9), Inches(4.55), Inches(2.45), Inches(1.16), "Accuracy", pct(ver4["multiclass"]["accuracy"]), TEAL)
    add_metric(slide, Inches(3.55), Inches(4.55), Inches(2.45), Inches(1.16), "Macro F1", pct(ver4["multiclass"]["macro_avg"]["f1"]), GREEN)
    add_metric(slide, Inches(6.2), Inches(4.55), Inches(2.45), Inches(1.16), "Mismatches", str(len(ver4["mismatches"])), ORANGE)
    add_text(slide, Inches(0.88), Inches(6.38), Inches(4.3), Inches(0.26), "2026-04-29", 12, RGBColor(214, 224, 232))

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "1. 課題と選んだ理由", "誤情報を早く見つけ、根拠付きで確認する支援が必要")
    add_text(slide, Inches(0.75), Inches(1.45), Inches(5.4), Inches(0.62), "取り組む課題", 24, NAVY, True)
    add_bullets(
        slide,
        Inches(0.78),
        Inches(2.15),
        Inches(5.25),
        Inches(2.25),
        [
            "SNSやニュース記事では、誤情報が短時間で広がりやすい",
            "人手で全件確認するには時間がかかり、優先順位づけが難しい",
            "単に真偽だけでなく、なぜそう判定したかの説明が必要",
        ],
        16,
    )
    add_panel(slide, Inches(6.55), Inches(1.48), Inches(5.82), Inches(3.15), BLUE_SOFT)
    add_text(slide, Inches(6.88), Inches(1.82), Inches(5.15), Inches(0.34), "選んだ理由", 21, TEAL, True)
    add_text(
        slide,
        Inches(6.9),
        Inches(2.38),
        Inches(4.95),
        Inches(1.25),
        "LLMは文章理解が得意だが、真偽判定には根拠確認が必要。そこで、LLMと外部根拠確認を組み合わせた実用的な確認支援ツールを検証した。",
        20,
        INK,
    )
    add_text(slide, Inches(0.82), Inches(5.35), Inches(11.3), Inches(0.52), "目指したもの", 17, ORANGE, True)
    add_text(
        slide,
        Inches(0.82),
        Inches(5.9),
        Inches(11.6),
        Inches(0.64),
        "大量の情報を一次スクリーニングし、人が確認すべき投稿を根拠付きで示す。",
        23,
        NAVY,
        True,
    )
    add_footer(slide, 2)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "2. 使用したデータについて", "100件の実記事・主張を5段階ラベルで評価")
    add_metric(slide, Inches(0.82), Inches(1.52), Inches(2.2), Inches(1.18), "Total", str(len(dataset["cases"])), NAVY, "評価対象")
    add_metric(slide, Inches(3.32), Inches(1.52), Inches(2.2), Inches(1.18), "Labels", str(len(verdict_counts)), TEAL, "5段階分類")
    add_metric(slide, Inches(5.82), Inches(1.52), Inches(2.2), Inches(1.18), "Each Label", str(min(verdict_counts.values())), GREEN, "各20件")
    add_metric(slide, Inches(8.32), Inches(1.52), Inches(2.2), Inches(1.18), "Skipped", str(ver4["meta"]["skipped_count"]), ORANGE, "実行時")
    add_panel(slide, Inches(0.82), Inches(3.15), Inches(5.55), Inches(2.6), WHITE)
    add_text(slide, Inches(1.1), Inches(3.42), Inches(4.8), Inches(0.28), "ラベル構成", 18, NAVY, True)
    label_order = ["正確", "ほぼ正確", "判断保留", "不正確", "誤り"]
    for i, label in enumerate(label_order):
        add_bar(slide, Inches(1.12), Inches(3.92 + i * 0.32), Inches(2.9), Inches(0.11), label, verdict_counts[label] / 20, TEAL)
        add_text(slide, Inches(5.2), Inches(3.86 + i * 0.32), Inches(0.5), Inches(0.16), f"{verdict_counts[label]}件", 8, SUBTLE, False, PP_ALIGN.RIGHT)
    add_panel(slide, Inches(6.82), Inches(3.15), Inches(5.55), Inches(2.6), WHITE)
    add_text(slide, Inches(7.1), Inches(3.42), Inches(4.8), Inches(0.28), "ドメイン構成", 18, NAVY, True)
    domain_rows = domain_counts.most_common(5)
    max_domain = max(count for _, count in domain_rows)
    for i, (domain, count) in enumerate(domain_rows):
        add_bar(slide, Inches(7.12), Inches(3.92 + i * 0.32), Inches(2.9), Inches(0.11), str(domain), count / max_domain, ORANGE)
        add_text(slide, Inches(11.2), Inches(3.86 + i * 0.32), Inches(0.5), Inches(0.16), f"{count}件", 8, SUBTLE, False, PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.9), Inches(6.25), Inches(11.45), Inches(0.3), "評価しやすいようにラベル数を揃え、正確性だけでなくラベルごとの弱点も見えるようにした。", 16, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 3)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "3. 初期仮説", "文章だけでは、ニュースの真偽判定は安定しない")
    add_text(slide, Inches(0.75), Inches(1.45), Inches(5.4), Inches(0.62), "なぜ難しいのか", 24, NAVY, True)
    add_bullets(
        slide,
        Inches(0.78),
        Inches(2.15),
        Inches(5.25),
        Inches(2.25),
        [
            "数字・時期・固有名詞の小さなズレが判定を左右する",
            "一部だけ正しい主張は、単純な真偽分類では扱いにくい",
            "根拠不足や学説未確立の話題は、断定すると誤判定になりやすい",
        ],
        16,
    )
    add_panel(slide, Inches(6.55), Inches(1.48), Inches(5.82), Inches(3.15), BLUE_SOFT)
    add_text(slide, Inches(6.88), Inches(1.82), Inches(5.15), Inches(0.34), "Ver4の考え方", 21, TEAL, True)
    add_text(
        slide,
        Inches(6.9),
        Inches(2.38),
        Inches(4.95),
        Inches(1.25),
        "文章の印象だけで決めず、外部根拠と照合してから、正確・ほぼ正確・判断保留・不正確・誤りへ整理する。",
        20,
        INK,
    )
    add_text(slide, Inches(0.82), Inches(5.35), Inches(11.3), Inches(0.52), "発表での一言", 17, ORANGE, True)
    add_text(
        slide,
        Inches(0.82),
        Inches(5.9),
        Inches(11.6),
        Inches(0.64),
        "Ver4は、LLMの文脈理解と外部根拠確認を組み合わせ、断定しすぎない5段階判定を目指した。",
        23,
        NAVY,
        True,
    )
    add_footer(slide, 4)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "3. 分析アプローチ", "入力から最終ラベルまでを段階的に確認する")
    steps = [
        ("入力", "記事本文・主張・確認対象を受け取る", TEAL),
        ("ローカル一次判定", "注意度や候補ラベルを先に整理", GREEN),
        ("GPT一次レビュー", "文脈理解と論点抽出を行う", GOLD),
        ("Gemini根拠確認", "外部情報と主張を照合する", ORANGE),
        ("補正ルール", "5段階ラベルへキャリブレーション", RED),
        ("出力", "判定・理由・確認リンクを返す", NAVY),
    ]
    x0 = Inches(0.55)
    y0 = Inches(2.0)
    w = Inches(1.72)
    h = Inches(1.25)
    gap = Inches(0.28)
    for i, (title, body, color) in enumerate(steps):
        left = x0 + i * (w + gap)
        add_step(slide, left, y0, w, h, title, body, color, str(i + 1))
        if i < len(steps) - 1:
            add_arrow(slide, left + w + Inches(0.04), y0 + Inches(0.48), Inches(0.22), Inches(0.25), SUBTLE)
    add_panel(slide, Inches(1.05), Inches(4.45), Inches(11.25), Inches(1.42), MINT)
    add_text(slide, Inches(1.38), Inches(4.78), Inches(10.6), Inches(0.36), "重要なのは、LLMの出力をそのまま採用しないこと", 23, NAVY, True, PP_ALIGN.CENTER)
    add_text(
        slide,
        Inches(1.45),
        Inches(5.34),
        Inches(10.35),
        Inches(0.3),
        "外部根拠との対応関係を見て、最終判定をルールで補正する。",
        15,
        SUBTLE,
        False,
        PP_ALIGN.CENTER,
    )
    add_footer(slide, 5)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "3. 5段階判定", "現実の主張にあるグラデーションを表現する")
    labels = [
        ("正確", "主要部分が根拠と一致", GREEN, GREEN_SOFT),
        ("ほぼ正確", "大筋は正しいが軽微なズレ", TEAL, MINT),
        ("判断保留", "根拠不足・未確定で断定不可", GOLD, RGBColor(251, 246, 226)),
        ("不正確", "一部真実だが重要部分が誤り", ORANGE, ORANGE_SOFT),
        ("誤り", "核心部分が明確に否定される", RED, RED_SOFT),
    ]
    for i, (label, desc, color, fill) in enumerate(labels):
        top = Inches(1.55 + i * 0.86)
        add_panel(slide, Inches(0.8), top, Inches(11.8), Inches(0.62), fill, line=RGBColor(220, 226, 228))
        add_text(slide, Inches(1.08), top + Inches(0.14), Inches(1.75), Inches(0.22), label, 17, color, True)
        add_text(slide, Inches(3.0), top + Inches(0.14), Inches(8.85), Inches(0.22), desc, 16, INK)
    add_text(slide, Inches(0.9), Inches(6.25), Inches(11.4), Inches(0.4), "ポイント: 「完全に正しい / 完全に誤り」だけでなく、中間ケースを明示的に扱う。", 18, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 6)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "4. 分析結果", "100件のbalanced datasetで5段階分類を評価")
    add_metric(slide, Inches(0.78), Inches(1.52), Inches(2.15), Inches(1.18), "Sample", str(ver4["meta"]["sample_count"]), NAVY, "5ラベル各20件")
    add_metric(slide, Inches(3.18), Inches(1.52), Inches(2.15), Inches(1.18), "Accuracy", pct(ver4["multiclass"]["accuracy"]), TEAL)
    add_metric(slide, Inches(5.58), Inches(1.52), Inches(2.15), Inches(1.18), "Macro F1", pct(ver4["multiclass"]["macro_avg"]["f1"]), GREEN)
    add_metric(slide, Inches(7.98), Inches(1.52), Inches(2.15), Inches(1.18), "Binary F1", pct(ver4["binary_fake_positive"]["f1"]), ORANGE, "誤り検出")
    add_metric(slide, Inches(10.38), Inches(1.52), Inches(2.15), Inches(1.18), "Skipped", str(ver4["meta"]["skipped_count"]), RED)
    add_picture_fit(slide, VER4_PLOTS_DIR / "per_class_metrics.png", Inches(0.75), Inches(3.18), Inches(5.7), Inches(2.78))
    add_picture_fit(slide, VER4_PLOTS_DIR / "summary_metrics.png", Inches(6.78), Inches(3.18), Inches(5.7), Inches(2.78))
    add_footer(slide, 7)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "4. 考察", "LLM出力を5段階ラベルへどう写像するかを改善")
    rules = [
        ("核心が正しい + 軽微差分", "数字・時期・表現差は、必要に応じて「ほぼ正確」へ", GREEN, GREEN_SOFT),
        ("一部真実 + 重要な誤り", "正しい要素があっても、結論を歪める場合は「不正確」へ", ORANGE, ORANGE_SOFT),
        ("存在しない発言・明確な反証", "根拠で否定できる場合は「誤り」へ", RED, RED_SOFT),
        ("論争中・根拠不足", "歴史論争や未確定情報は「判断保留」へ", GOLD, RGBColor(251, 246, 226)),
    ]
    for i, (title, body, color, fill) in enumerate(rules):
        left = Inches(0.78 + (i % 2) * 6.05)
        top = Inches(1.58 + (i // 2) * 1.82)
        add_panel(slide, left, top, Inches(5.55), Inches(1.38), fill)
        add_text(slide, left + Inches(0.22), top + Inches(0.22), Inches(5.05), Inches(0.28), title, 17, color, True)
        add_text(slide, left + Inches(0.22), top + Inches(0.68), Inches(5.05), Inches(0.45), body, 13, INK)
    add_panel(slide, Inches(1.18), Inches(5.62), Inches(10.95), Inches(0.72), BLUE_SOFT)
    add_text(
        slide,
        Inches(1.35),
        Inches(5.85),
        Inches(10.55),
        Inches(0.24),
        "モデル自体を再学習したのではなく、判定ロジックのキャリブレーションを行った。",
        18,
        NAVY,
        True,
        PP_ALIGN.CENTER,
    )
    add_footer(slide, 8)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "4. 残った誤分類", "正確 / ほぼ正確 / 不正確の境界が主な残課題")
    add_picture_fit(slide, VER4_PLOTS_DIR / "confusion_matrix.png", Inches(0.72), Inches(1.42), Inches(5.6), Inches(4.58))
    add_text(slide, Inches(6.75), Inches(1.45), Inches(4.6), Inches(0.34), f"残ミスマッチ: {len(ver4['mismatches'])}件", 22, NAVY, True)
    y = 2.06
    for (truth, pred), count in mismatch_patterns[:5]:
        add_panel(slide, Inches(6.78), Inches(y), Inches(5.15), Inches(0.54), WHITE)
        add_text(slide, Inches(7.05), Inches(y + 0.13), Inches(3.2), Inches(0.18), f"{truth} → {pred}", 13, INK, True)
        add_text(slide, Inches(11.0), Inches(y + 0.13), Inches(0.62), Inches(0.18), f"{count}件", 13, ORANGE, True, PP_ALIGN.RIGHT)
        y += 0.72
    add_text(
        slide,
        Inches(6.8),
        Inches(5.88),
        Inches(5.25),
        Inches(0.42),
        "今後は固有名詞パッチではなく、境界判定の構造ルールとして改善する。",
        15,
        SUBTLE,
    )
    add_footer(slide, 9)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "4. バージョン比較", "保存済み評価ではVer4が最も高い")
    versions = [
        ("Ver2", ver2["multiclass"]["accuracy"], ver2["multiclass"]["macro_avg"]["f1"], BLUE_SOFT, NAVY),
        ("Ver3", ver3["multiclass"]["accuracy"], ver3["multiclass"]["macro_avg"]["f1"], MINT, TEAL),
        ("Ver4", ver4["multiclass"]["accuracy"], ver4["multiclass"]["macro_avg"]["f1"], GREEN_SOFT, GREEN),
    ]
    for i, (name, acc, f1, fill, accent_color) in enumerate(versions):
        top = Inches(1.55 + i * 1.03)
        add_panel(slide, Inches(0.95), top, Inches(11.3), Inches(0.72), fill)
        add_text(slide, Inches(1.2), top + Inches(0.22), Inches(1.0), Inches(0.18), name, 16, accent_color, True)
        add_bar(slide, Inches(2.45), top + Inches(0.19), Inches(3.55), Inches(0.16), "Accuracy", acc, accent_color)
        add_bar(slide, Inches(7.0), top + Inches(0.19), Inches(3.55), Inches(0.16), "Macro F1", f1, accent_color)
    add_text(
        slide,
        Inches(0.95),
        Inches(6.12),
        Inches(11.35),
        Inches(0.34),
        "注意: 既存100件での評価結果。未知データへの汎化性能はblind testで確認する必要がある。",
        15,
        RED,
        True,
        PP_ALIGN.CENTER,
    )
    add_footer(slide, 10)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "5. ビジネスインパクト", "人の確認作業を置き換えるのではなく、優先順位づけを支援する")
    impacts = [
        ("一次スクリーニング", "大量の投稿・記事から確認優先度の高いものを抽出", TEAL, MINT),
        ("調査時間の短縮", "主張・理由・確認リンクをまとめ、ファクトチェックの初動を速くする", GREEN, GREEN_SOFT),
        ("説明責任の補助", "判定だけでなく、なぜそう見たかを文章で残せる", GOLD, RGBColor(251, 246, 226)),
        ("運用リスク低減", "判断保留を残すことで、根拠不足の断定を避ける", ORANGE, ORANGE_SOFT),
    ]
    for i, (title, body, color, fill) in enumerate(impacts):
        left = Inches(0.78 + (i % 2) * 6.05)
        top = Inches(1.55 + (i // 2) * 1.78)
        add_panel(slide, left, top, Inches(5.55), Inches(1.32), fill)
        add_text(slide, left + Inches(0.22), top + Inches(0.22), Inches(5.05), Inches(0.25), title, 17, color, True)
        add_text(slide, left + Inches(0.22), top + Inches(0.65), Inches(5.0), Inches(0.42), body, 13, INK)
    add_panel(slide, Inches(1.12), Inches(5.58), Inches(11.05), Inches(0.82), BLUE_SOFT)
    add_text(slide, Inches(1.35), Inches(5.86), Inches(10.55), Inches(0.25), "活用場面: メディア監視、広報・危機管理、自治体や企業の情報確認フロー", 17, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 11)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "6. 工夫点", "精度だけでなく、説明しやすさと改善しやすさを重視")
    ideas = [
        ("評価bundle化", "JSON・CSV・XLSX・可視化を保存し、結果をあとから検証できる"),
        ("5段階ラベル", "完全な正誤だけでなく、軽微差分と判断保留を扱える"),
        ("エラー分析", "ミスマッチを境界パターンで分類し、改善対象を明確化"),
        ("キャリブレーション", "モデル再学習ではなく、ラベルマッピングを改善"),
    ]
    for i, (title, body) in enumerate(ideas):
        add_step(slide, Inches(1.02), Inches(1.45 + i * 1.08), Inches(11.2), Inches(0.76), title, body, [TEAL, GREEN, GOLD, ORANGE][i], str(i + 1))
    add_footer(slide, 12)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "6. 今後さらに実施したいこと", "既存100件をdev、新規100件をblind testとして使う")
    roadmap = [
        ("1", "既存100件をdevセット化", "現在の改善効果を継続確認する", TEAL),
        ("2", "新規blind 100件を作成", "未知データでの実力を測る", GREEN),
        ("3", "Ver2/3/4/5を同条件で評価", "方式差を公平に比較する", GOLD),
        ("4", "構造ルールで改善", "固有名詞単位の過学習を避ける", ORANGE),
        ("5", "devとblindの両方で確認", "既存性能を落とさず一般化を狙う", RED),
    ]
    for i, (num, title, body, color) in enumerate(roadmap):
        top = Inches(1.45 + i * 0.95)
        add_step(slide, Inches(1.05), top, Inches(11.1), Inches(0.64), title, body, color, num)
    add_footer(slide, 13)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "まとめ", "Ver4は「確認支援ツール」としての実用性を高めた構成")
    add_panel(slide, Inches(1.05), Inches(1.58), Inches(11.1), Inches(1.2), MINT)
    add_text(
        slide,
        Inches(1.32),
        Inches(1.92),
        Inches(10.58),
        Inches(0.34),
        "既存LLM + 外部根拠確認 + ルールベース補正による真偽判定パイプライン",
        22,
        NAVY,
        True,
        PP_ALIGN.CENTER,
    )
    add_bullets(
        slide,
        Inches(1.45),
        Inches(3.35),
        Inches(10.35),
        Inches(1.95),
        [
            "GPTで文脈と論点を整理し、Geminiで外部根拠を確認する",
            "真偽を二値ではなく5段階で扱い、中間ケースを表現する",
            "改善の中心は、モデル再学習ではなく判定ロジックのキャリブレーション",
            "次の課題は、新規blind testによる汎化性能の検証",
        ],
        17,
    )
    add_text(slide, Inches(1.25), Inches(6.16), Inches(10.75), Inches(0.38), "発表では「文章だけでなく根拠と照合して判定する」点を強調する。", 18, TEAL, True, PP_ALIGN.CENTER)
    add_footer(slide, 14)

    prs.save(OUTPUT_PPTX)
    OUTPUT_NOTES.write_text(build_notes(ver4, dataset, mismatch_patterns), encoding="utf-8")


if __name__ == "__main__":
    build_presentation()
    print(f"saved: {OUTPUT_PPTX}")
    print(f"saved: {OUTPUT_NOTES}")
