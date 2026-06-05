from __future__ import annotations

import csv
import json
from collections import Counter
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PRESENTATION_DIR = Path(__file__).resolve().parent
PROJECT_ROOT = PRESENTATION_DIR.parent

SOURCE_PPTX = PRESENTATION_DIR / "fake_news_checker_presentation_20260401.pptx"
DATASET_PATH = PROJECT_ROOT / "testdata" / "shared" / "real_article_dataset_v2.json"
EVAL_DIR = PROJECT_ROOT / "Ver2" / "evaluation_outputs" / "20260411-0509"
EVAL_PATH = EVAL_DIR / "eval_real_article_dataset_v2_use_gemini.json"
PREDICTIONS_PATH = EVAL_DIR / "predictions_real_article_dataset_v2_use_gemini.json"
CSV_PATH = EVAL_DIR / "Ver2_real_article_dataset_v2_with_predicted_verdict_attention_score.csv"
PLOTS_DIR = EVAL_DIR / "plots"

OUTPUT_PPTX = PRESENTATION_DIR / "fake_news_checker_ver2_eval_summary_20260411.pptx"
OUTPUT_NOTES = PRESENTATION_DIR / "fake_news_checker_ver2_eval_summary_20260411_notes.md"

BG = RGBColor(247, 244, 238)
INK = RGBColor(31, 42, 52)
SUBTLE = RGBColor(98, 107, 118)
LINE = RGBColor(217, 211, 202)
WHITE = RGBColor(255, 255, 255)
NAVY = RGBColor(22, 39, 57)
TEAL = RGBColor(19, 118, 122)
ORANGE = RGBColor(201, 109, 49)
GREEN = RGBColor(69, 142, 95)
GOLD = RGBColor(192, 145, 43)
RED = RGBColor(181, 61, 61)
MINT = RGBColor(234, 243, 241)
SAND = RGBColor(243, 238, 229)

FONT_TITLE = "Yu Gothic"
FONT_BODY = "Yu Gothic"


@dataclass
class EvalSummary:
    analysis_datetime: str
    case_count: int
    accuracy: float
    macro_precision: float
    macro_recall: float
    macro_f1: float
    weighted_f1: float
    binary_precision: float
    binary_recall: float
    binary_f1: float
    per_class: dict[str, dict]
    mismatch_count: int
    mismatch_patterns: list[tuple[tuple[str, str], int]]


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def load_csv_rows(path: Path) -> list[dict[str, str]]:
    with path.open(encoding="utf-8-sig", newline="") as f:
        return list(csv.DictReader(f))


def build_eval_summary() -> EvalSummary:
    report = load_json(EVAL_PATH)
    predictions = load_json(PREDICTIONS_PATH)
    mismatch_counts = Counter((item["truth"], item["pred"]) for item in report["mismatches"])
    return EvalSummary(
        analysis_datetime=str(predictions["meta"].get("analysis_datetime") or ""),
        case_count=int(predictions["meta"].get("case_count") or 0),
        accuracy=float(report["multiclass"]["accuracy"]),
        macro_precision=float(report["multiclass"]["macro_avg"]["precision"]),
        macro_recall=float(report["multiclass"]["macro_avg"]["recall"]),
        macro_f1=float(report["multiclass"]["macro_avg"]["f1"]),
        weighted_f1=float(report["multiclass"]["weighted_avg"]["f1"]),
        binary_precision=float(report["binary_fake_positive"]["precision"]),
        binary_recall=float(report["binary_fake_positive"]["recall"]),
        binary_f1=float(report["binary_fake_positive"]["f1"]),
        per_class=dict(report["multiclass"]["per_class"]),
        mismatch_count=len(report["mismatches"]),
        mismatch_patterns=mismatch_counts.most_common(),
    )


def build_dataset_summary() -> tuple[Counter, Counter]:
    payload = load_json(DATASET_PATH)
    verdict_counts = Counter(case.get("expected_verdict") for case in payload["cases"])
    domain_counts = Counter(case.get("expected_domain") for case in payload["cases"])
    return verdict_counts, domain_counts


def build_representative_mismatches() -> list[dict[str, str]]:
    rows_by_id = {row["id"]: row for row in load_csv_rows(CSV_PATH)}
    report = load_json(EVAL_PATH)
    targets = [
        ("正確", "ほぼ正確"),
        ("ほぼ正確", "正確"),
        ("ほぼ正確", "不正確"),
        ("誤り", "不正確"),
    ]
    picks: list[dict[str, str]] = []
    for truth, pred in targets:
        for item in report["mismatches"]:
            if item["truth"] == truth and item["pred"] == pred:
                row = rows_by_id[item["id"]]
                picks.append(
                    {
                        "flow": f"{truth} → {pred}",
                        "id": item["id"],
                        "claim": row["analysis_text"].strip(),
                    }
                )
                break
    return picks


def pct(value: float) -> str:
    return f"{value * 100:.1f}%"


def short_text(text: str, limit: int = 54) -> str:
    if len(text) <= limit:
        return text
    return text[: limit - 1] + "…"


def set_slide_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_panel(slide, left, top, width, height, fill_rgb=WHITE, line_rgb=LINE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = line_rgb
    shape.line.width = Pt(1.15)
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str = "",
    font_size: int = 20,
    color: RGBColor = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if text:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = FONT_BODY
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        p.alignment = align
    return box


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    items: list[str],
    font_size: int = 18,
    color: RGBColor = INK,
    spacing: int = 6,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.space_after = Pt(spacing)
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = FONT_BODY
            run.font.size = Pt(font_size)
            run.font.color.rgb = color


def add_header(slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(slide, Inches(0.62), Inches(0.35), Inches(9.6), Inches(0.48), title, 25, INK, True)
    if subtitle:
        add_textbox(slide, Inches(0.64), Inches(0.84), Inches(10.4), Inches(0.28), subtitle, 11, SUBTLE, False)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.63), Inches(1.12), Inches(1.55), Inches(0.07))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def add_footer(slide, label: str) -> None:
    add_textbox(slide, Inches(11.75), Inches(7.0), Inches(1.0), Inches(0.2), label, 10, SUBTLE, False, PP_ALIGN.RIGHT)


def add_picture_contained(slide, image_path: Path, left, top, width, height) -> None:
    pic = slide.shapes.add_picture(str(image_path), left, top)
    scale = min(width / pic.width, height / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = left + int((width - pic.width) / 2)
    pic.top = top + int((height - pic.height) / 2)


def add_metric_card(slide, left, top, width, title: str, value: str, fill_rgb=MINT) -> None:
    add_panel(slide, left, top, width, Inches(1.05), fill_rgb=fill_rgb)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.14), width - Inches(0.36), Inches(0.25), title, 12, SUBTLE, False)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.38), width - Inches(0.36), Inches(0.42), value, 24, INK, True)


def build_presentation() -> tuple[Presentation, list[str]]:
    summary = build_eval_summary()
    verdict_counts, domain_counts = build_dataset_summary()
    examples = build_representative_mismatches()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    notes: list[str] = []

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_panel(slide, Inches(0.5), Inches(0.45), Inches(12.33), Inches(6.45), fill_rgb=WHITE)
    add_textbox(slide, Inches(0.95), Inches(1.05), Inches(7.6), Inches(0.72), "フェイクニュースチェッカー", 28, INK, True)
    add_textbox(slide, Inches(0.95), Inches(1.72), Inches(8.6), Inches(0.7), "Ver2 最新評価サマリー", 24, TEAL, True)
    add_textbox(
        slide,
        Inches(0.98),
        Inches(2.5),
        Inches(6.7),
        Inches(1.3),
        "アップロード済みの発表資料、評価 CSV、5 種類の可視化画像をもとに、\n最新 100 件 rerun の結果を要約。",
        18,
        SUBTLE,
        False,
    )
    add_metric_card(slide, Inches(8.0), Inches(1.2), Inches(1.8), "Accuracy", pct(summary.accuracy), SAND)
    add_metric_card(slide, Inches(9.95), Inches(1.2), Inches(1.8), "Macro F1", f"{summary.macro_f1:.3f}", MINT)
    add_metric_card(slide, Inches(8.0), Inches(2.55), Inches(1.8), "Binary F1", f"{summary.binary_f1:.3f}", MINT)
    add_metric_card(slide, Inches(9.95), Inches(2.55), Inches(1.8), "Mismatches", str(summary.mismatch_count), SAND)
    add_textbox(slide, Inches(0.98), Inches(5.75), Inches(7.6), Inches(0.32), f"参照元: {SOURCE_PPTX.name}", 11, SUBTLE)
    add_textbox(slide, Inches(0.98), Inches(6.05), Inches(7.6), Inches(0.32), f"評価日時: {summary.analysis_datetime}", 11, SUBTLE)
    add_footer(slide, "1")
    notes.append("表紙。既存資料の続編として、Ver2 の最新評価を短く説明する。Accuracy 84.0%、Macro F1 0.837、Binary F1 0.889 を最初に示す。")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "1. 目的と評価設定", "既存発表資料の問題意識を引き継ぎつつ、今回は Ver2 の最新 rerun に絞って報告")
    add_panel(slide, Inches(0.6), Inches(1.45), Inches(5.95), Inches(5.45), fill_rgb=WHITE)
    add_panel(slide, Inches(6.75), Inches(1.45), Inches(5.95), Inches(5.45), fill_rgb=WHITE)
    add_textbox(slide, Inches(0.9), Inches(1.7), Inches(3.0), Inches(0.3), "プロジェクトの狙い", 18, INK, True)
    add_bullets(
        slide,
        Inches(0.88),
        Inches(2.08),
        Inches(5.1),
        Inches(2.8),
        [
            "SNS で流れる主張を 5 区分で整理し、投稿前の確認を支援する。",
            "完全自動の断定ではなく、説明つきの判定補助を目指す。",
            "ローカル判定と Gemini 外部根拠比較を組み合わせて境界ケースを扱う。",
        ],
        18,
    )
    add_textbox(slide, Inches(0.9), Inches(4.95), Inches(3.0), Inches(0.3), "今回の評価条件", 18, INK, True)
    add_bullets(
        slide,
        Inches(0.88),
        Inches(5.32),
        Inches(5.2),
        Inches(1.2),
        [
            "dataset: real_article_dataset_v2.json",
            "100 件、5 区分は各 20 件の balanced 構成",
            "Gemini 有効で rerun、bundle を保存",
        ],
        16,
    )
    add_textbox(slide, Inches(7.05), Inches(1.7), Inches(3.6), Inches(0.3), "データセット構成", 18, INK, True)
    verdict_line = " / ".join(f"{k} {v}件" for k, v in verdict_counts.items())
    domain_line = " / ".join(f"{k} {v}" for k, v in domain_counts.items())
    add_bullets(
        slide,
        Inches(7.02),
        Inches(2.08),
        Inches(5.2),
        Inches(1.8),
        [
            f"ラベル分布: {verdict_line}",
            f"ドメイン分布: {domain_line}",
            "attention_score つき CSV と plots を同時出力",
        ],
        16,
    )
    add_textbox(slide, Inches(7.05), Inches(4.2), Inches(3.6), Inches(0.3), "評価の見方", 18, INK, True)
    add_bullets(
        slide,
        Inches(7.02),
        Inches(4.58),
        Inches(5.15),
        Inches(1.9),
        [
            "正確 / ほぼ正確の境界がいちばん難しい。",
            "誤り recall を上げつつ、過検知を抑える必要がある。",
            "今回の資料では数値だけでなく、誤分類の流れも確認する。",
        ],
        16,
    )
    add_footer(slide, "2")
    notes.append("2枚目は目的と条件。100件 balanced dataset を Gemini 付きで再評価したこと、実運用では確認支援ツールとして位置づけていることを説明する。")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "2. 全体結果", "Accuracy / F1 の水準と、今回の rerun で見えている全体傾向")
    add_panel(slide, Inches(0.62), Inches(1.42), Inches(7.35), Inches(5.65), fill_rgb=WHITE)
    add_panel(slide, Inches(8.15), Inches(1.42), Inches(4.55), Inches(5.65), fill_rgb=WHITE)
    add_picture_contained(slide, PLOTS_DIR / "summary_metrics.png", Inches(0.8), Inches(1.72), Inches(6.95), Inches(4.95))
    add_textbox(slide, Inches(8.45), Inches(1.7), Inches(2.6), Inches(0.3), "主要指標", 18, INK, True)
    add_bullets(
        slide,
        Inches(8.4),
        Inches(2.08),
        Inches(3.9),
        Inches(2.0),
        [
            f"Accuracy: {pct(summary.accuracy)}",
            f"Macro F1: {summary.macro_f1:.3f}",
            f"Weighted F1: {summary.weighted_f1:.3f}",
            f"誤り binary F1: {summary.binary_f1:.3f}",
        ],
        18,
    )
    add_textbox(slide, Inches(8.45), Inches(4.45), Inches(2.6), Inches(0.3), "読み取り", 18, INK, True)
    add_bullets(
        slide,
        Inches(8.4),
        Inches(4.82),
        Inches(3.95),
        Inches(1.5),
        [
            "全体では 84% まで到達し、支援ツールとしてはかなり実用的。",
            "一方で境界ラベルの取り扱いにまだ粗さが残る。",
        ],
        16,
    )
    add_footer(slide, "3")
    notes.append("3枚目は全体結果。Accuracy 84.0%、Macro F1 0.837、誤り binary F1 0.889。以前より改善したが、境界ラベルの揺れが残る点を添える。")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "3. クラス別の特徴", "クラスごとに見ると、得意・苦手がかなり分かれている")
    add_panel(slide, Inches(0.62), Inches(1.42), Inches(8.05), Inches(5.65), fill_rgb=WHITE)
    add_panel(slide, Inches(8.9), Inches(1.42), Inches(3.8), Inches(5.65), fill_rgb=WHITE)
    add_picture_contained(slide, PLOTS_DIR / "per_class_metrics.png", Inches(0.8), Inches(1.72), Inches(7.65), Inches(4.95))
    add_textbox(slide, Inches(9.15), Inches(1.7), Inches(2.8), Inches(0.3), "ポイント", 18, INK, True)
    add_bullets(
        slide,
        Inches(9.12),
        Inches(2.08),
        Inches(3.2),
        Inches(3.8),
        [
            f"判断保留: F1 {summary.per_class['判断保留']['f1']:.3f}",
            f"不正確: recall {summary.per_class['不正確']['recall']:.3f}",
            f"ほぼ正確: recall {summary.per_class['ほぼ正確']['recall']:.3f}",
            f"誤り: recall {summary.per_class['誤り']['recall']:.3f}",
            "保留・不正確は強いが、ほぼ正確の回収率がまだ低い。",
        ],
        16,
    )
    add_footer(slide, "4")
    notes.append("4枚目はクラス別。判断保留と不正確は安定している一方、ほぼ正確 recall 0.600 がボトルネック。誤り recall 0.800 もまだ伸ばしたい。")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "4. 混同行列", "誤分類の向きを見ると、改善ポイントがかなり明確")
    add_panel(slide, Inches(0.62), Inches(1.42), Inches(8.0), Inches(5.65), fill_rgb=WHITE)
    add_panel(slide, Inches(8.85), Inches(1.42), Inches(3.85), Inches(5.65), fill_rgb=WHITE)
    add_picture_contained(slide, PLOTS_DIR / "confusion_matrix.png", Inches(0.8), Inches(1.72), Inches(7.6), Inches(4.95))
    add_textbox(slide, Inches(9.1), Inches(1.7), Inches(2.8), Inches(0.3), "主な流れ", 18, INK, True)
    flow_items = [f"{truth} → {pred}: {count}件" for (truth, pred), count in summary.mismatch_patterns]
    add_bullets(slide, Inches(9.08), Inches(2.08), Inches(3.2), Inches(2.3), flow_items, 16)
    add_textbox(slide, Inches(9.1), Inches(4.8), Inches(2.8), Inches(0.3), "解釈", 18, INK, True)
    add_bullets(
        slide,
        Inches(9.08),
        Inches(5.18),
        Inches(3.2),
        Inches(1.15),
        [
            "境界ケースをやや安全側に倒している。",
            "誤りを不正確に落とす 4 件が次の重点改善点。",
        ],
        15,
    )
    add_footer(slide, "5")
    notes.append("5枚目は混同行列。ほぼ正確→正確が6件、誤り→不正確が4件。境界を安全側に倒した結果だが、今後は誤り recall の改善が重要。")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "5. 代表的な誤分類", "ラベルごとに、どこで迷っているかを具体例で確認")
    y = Inches(1.55)
    colors = [MINT, SAND, WHITE, MINT]
    for idx, example in enumerate(examples):
        add_panel(slide, Inches(0.72), y + Inches(1.23 * idx), Inches(11.9), Inches(0.98), fill_rgb=colors[idx % len(colors)])
        add_textbox(slide, Inches(0.98), y + Inches(1.23 * idx) + Inches(0.12), Inches(2.2), Inches(0.22), example["flow"], 15, TEAL, True)
        add_textbox(slide, Inches(3.1), y + Inches(1.23 * idx) + Inches(0.11), Inches(8.95), Inches(0.56), short_text(example["claim"], 70), 17, INK, False)
        add_textbox(slide, Inches(0.98), y + Inches(1.23 * idx) + Inches(0.56), Inches(5.0), Inches(0.18), example["id"], 10, SUBTLE, False)
    add_textbox(slide, Inches(0.82), Inches(6.62), Inches(11.4), Inches(0.24), "特に『ほぼ正確』の境界と、『誤り』をどこで断定するかが改善の中心。", 15, SUBTLE, False)
    add_footer(slide, "6")
    notes.append("6枚目は代表誤分類。正確→ほぼ正確、ほぼ正確→正確、ほぼ正確→不正確、誤り→不正確を1件ずつ見せて、境界条件の調整が中心課題だと示す。")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "6. 改善方針", "今回の rerun から見えた、次の実装タスク")
    add_panel(slide, Inches(0.72), Inches(1.6), Inches(3.85), Inches(4.9), fill_rgb=WHITE)
    add_panel(slide, Inches(4.77), Inches(1.6), Inches(3.85), Inches(4.9), fill_rgb=WHITE)
    add_panel(slide, Inches(8.82), Inches(1.6), Inches(3.8), Inches(4.9), fill_rgb=WHITE)
    add_textbox(slide, Inches(1.02), Inches(1.95), Inches(2.6), Inches(0.3), "短期", 20, ORANGE, True)
    add_bullets(
        slide,
        Inches(0.98),
        Inches(2.35),
        Inches(3.1),
        Inches(2.9),
        [
            "誤り → 不正確 の 4 件を優先して潰す。",
            "陰謀論・歴史否定・明確な虚偽の強化ルールを点検する。",
            "100 件 rerun を回しながら小刻みに回帰確認する。",
        ],
        16,
    )
    add_textbox(slide, Inches(5.07), Inches(1.95), Inches(2.6), Inches(0.3), "中期", 20, TEAL, True)
    add_bullets(
        slide,
        Inches(5.03),
        Inches(2.35),
        Inches(3.1),
        Inches(2.9),
        [
            "ほぼ正確のルールをもう一段整理する。",
            "数値の近似、主導と共同関与、途中段階の省略を型で管理する。",
            "attention_score と判定ラベルの関係も再確認する。",
        ],
        16,
    )
    add_textbox(slide, Inches(9.1), Inches(1.95), Inches(2.6), Inches(0.3), "位置づけ", 20, GREEN, True)
    add_bullets(
        slide,
        Inches(9.06),
        Inches(2.35),
        Inches(3.05),
        Inches(2.9),
        [
            "現時点でも確認支援ツールとしては十分に有望。",
            "自動断定より、説明付きの一次仕分けに強みがある。",
            "発表では『使えるが、境界設計がまだ研究課題』とまとめる。",
        ],
        16,
    )
    add_footer(slide, "7")
    notes.append("7枚目は改善方針。短期は誤り recall、中期はほぼ正確の境界設計。最終的な位置づけは自動断定ではなく説明付き確認支援ツールとする。")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "参考: ダッシュボード", "アップロードされた評価画像をそのまま添付")
    add_panel(slide, Inches(0.62), Inches(1.42), Inches(12.0), Inches(5.65), fill_rgb=WHITE)
    add_picture_contained(slide, PLOTS_DIR / "evaluation_dashboard.png", Inches(0.78), Inches(1.68), Inches(11.65), Inches(5.15))
    add_footer(slide, "8")
    notes.append("最後はダッシュボード全体図。質疑ではこのスライドを開いて、混同行列・summary・per-class を一枚で見せられるようにする。")

    return prs, notes


def save_notes(notes: list[str]) -> None:
    lines = [
        "# Fake News Checker Ver2 Evaluation Summary Notes",
        "",
        f"- source_pptx: {SOURCE_PPTX.name}",
        f"- dataset: {DATASET_PATH.relative_to(PROJECT_ROOT)}",
        f"- evaluation: {EVAL_PATH.relative_to(PROJECT_ROOT)}",
        f"- predictions_csv: {CSV_PATH.relative_to(PROJECT_ROOT)}",
        "",
    ]
    for i, note in enumerate(notes, start=1):
        lines.append(f"## Slide {i}")
        lines.append(f"- {note}")
        lines.append("")
    OUTPUT_NOTES.write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    prs, notes = build_presentation()
    prs.save(OUTPUT_PPTX)
    save_notes(notes)
    print(OUTPUT_PPTX)
    print(OUTPUT_NOTES)


if __name__ == "__main__":
    main()
