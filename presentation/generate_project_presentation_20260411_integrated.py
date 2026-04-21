from __future__ import annotations

import json
from collections import Counter
from pathlib import Path

from generate_project_presentation import (
    BG,
    FONT_BODY,
    GREEN,
    INK,
    LINE,
    MINT,
    NAVY,
    ORANGE,
    PROJECT_ROOT,
    RED,
    SAND,
    SUBTLE,
    TEAL,
    WHITE,
    DatasetSummary,
    EvalSnapshot,
    add_bullets,
    add_connector,
    add_footer,
    add_header,
    add_labeled_image,
    add_metric_card,
    add_panel,
    add_process_box,
    add_textbox,
    build_dataset_summary,
    build_eval_snapshot,
    plot_path,
    score,
    set_slide_bg,
)
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PRESENTATION_DIR = Path(__file__).resolve().parent
DATASET_PATH = PROJECT_ROOT / "testdata" / "shared" / "real_article_dataset_v2.json"

LEGACY_VER2_EVAL_PATH = PROJECT_ROOT / "Ver2" / "evaluation_outputs" / "20260401-1213" / "eval_real_v2_use_gemini_after_rule_tune4.json"
LEGACY_VER2_PREDICTIONS_PATH = PROJECT_ROOT / "Ver2" / "evaluation_outputs" / "20260401-1213" / "predictions_real_v2_use_gemini_after_rule_tune4.json"
LEGACY_VER2_PLOTS_DIR = PROJECT_ROOT / "Ver2" / "evaluation_outputs" / "20260401-1213" / "plots"

CURRENT_VER2_EVAL_PATH = PROJECT_ROOT / "Ver2" / "evaluation_outputs" / "20260411-0509" / "eval_real_article_dataset_v2_use_gemini.json"
CURRENT_VER2_PREDICTIONS_PATH = PROJECT_ROOT / "Ver2" / "evaluation_outputs" / "20260411-0509" / "predictions_real_article_dataset_v2_use_gemini.json"
CURRENT_VER2_CSV_PATH = PROJECT_ROOT / "Ver2" / "evaluation_outputs" / "20260411-0509" / "Ver2_real_article_dataset_v2_with_predicted_verdict_attention_score.csv"
CURRENT_VER2_PLOTS_DIR = PROJECT_ROOT / "Ver2" / "evaluation_outputs" / "20260411-0509" / "plots"
SCREENSHOT_DIR = Path(r"c:\Users\oneuk\OneDrive\Desktop\画像素材\新しいフォルダー")
SCREENSHOT_INPUT = SCREENSHOT_DIR / "スクリーンショット (439).png"
SCREENSHOT_LOADING = SCREENSHOT_DIR / "スクリーンショット (464).png"
SCREENSHOT_RESULT = SCREENSHOT_DIR / "スクリーンショット (465).png"
SCREENSHOT_REASON = SCREENSHOT_DIR / "スクリーンショット 2026-04-11 065005.png"

OUTPUT_PPTX = PRESENTATION_DIR / "fake_news_checker_presentation_20260411_integrated.pptx"
OUTPUT_NOTES = PRESENTATION_DIR / "fake_news_checker_presentation_20260411_integrated_notes.md"

FOOTER_NOTE = "Fake News Checker | 2026-04-11 最新評価を本編へ統合"


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def pct(value: float) -> str:
    return f"{value * 100:.1f}%"


def add_cover_footer(slide) -> None:
    add_textbox(slide, Inches(1.15), Inches(6.08), Inches(5.2), Inches(0.28), "2025年7月期　中村隆之", 16, WHITE, False)


def add_picture_contained(slide, image_path: Path, left, top, width, height) -> None:
    pic = slide.shapes.add_picture(str(image_path), left, top)
    scale = min(width / pic.width, height / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = left + int((width - pic.width) / 2)
    pic.top = top + int((height - pic.height) / 2)


def add_screenshot_card(slide, image_path: Path, left, top, width, height, label: str, accent) -> None:
    frame = add_panel(slide, left, top, width, height, fill_rgb=WHITE, line_rgb=LINE)
    frame.line.color.rgb = LINE
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.09))
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()
    add_picture_contained(slide, image_path, left + Inches(0.12), top + Inches(0.18), width - Inches(0.24), height - Inches(0.62))
    add_textbox(slide, left + Inches(0.12), top + height - Inches(0.3), width - Inches(0.24), Inches(0.18), label, 12, SUBTLE, True, PP_ALIGN.CENTER)


def representative_mismatch_examples(csv_path: Path, eval_path: Path) -> list[str]:
    report = load_json(eval_path)
    csv_rows: dict[str, dict[str, str]] = {}
    with csv_path.open(encoding="utf-8-sig", newline="") as f:
        import csv

        for row in csv.DictReader(f):
            csv_rows[row["id"]] = row

    targets = [
        ("正確", "ほぼ正確"),
        ("ほぼ正確", "正確"),
        ("ほぼ正確", "不正確"),
        ("誤り", "不正確"),
    ]
    examples: list[str] = []
    for truth, pred in targets:
        for item in report["mismatches"]:
            if item["truth"] == truth and item["pred"] == pred:
                claim = csv_rows[item["id"]]["analysis_text"].strip()
                if len(claim) > 38:
                    claim = claim[:37] + "…"
                examples.append(f"{truth} -> {pred}: {claim}")
                break
    return examples


def build_notes(dataset: DatasetSummary, legacy: EvalSnapshot, current: EvalSnapshot) -> str:
    top_current = [f"{truth} -> {pred}: {count}件" for (truth, pred), count in current.mismatch_counts[:4]]
    lines = [
        "# Fake News Checker Integrated Presentation Notes",
        "",
        "- generated_on: 2026-04-11",
        f"- dataset: {DATASET_PATH.relative_to(PROJECT_ROOT)}",
        f"- legacy_ver2_eval: {LEGACY_VER2_EVAL_PATH.relative_to(PROJECT_ROOT)}",
        f"- current_ver2_eval: {CURRENT_VER2_EVAL_PATH.relative_to(PROJECT_ROOT)}",
        "",
        "## Slide 1",
        "- 20260401 本編資料の統合版であることを表紙で明示する。",
        "",
        "## Slide 2",
        "- 社会背景と個人的動機は元の本編資料を維持し、課題意識を先に共有する。",
        "",
        "## Slide 3",
        f"- dataset v2 は {dataset.case_count}件、5区分は各20件で balanced。",
        f"- ドメイン分布は 一般{dataset.domain_counts['一般']} / 医療{dataset.domain_counts['医療']} / 金融{dataset.domain_counts['金融']} / 災害{dataset.domain_counts['災害']} / 政治{dataset.domain_counts['政治']}。",
        "",
        "## Slide 4",
        "- 初期仮説・分析アプローチは元資料の構成を保ち、ローカル一次判定と Gemini 比較の流れを説明する。",
        "",
        "## Slide 5",
        "- UI の実行例として、入力 -> 判定中 -> 結果表示 の 3 ステップをスクリーンショットで示す。",
        "",
        "## Slide 6",
        "- 出力例では、要注意度・判定・理由がどのように返るかを説明する。",
        "",
        "## Slide 7",
        f"- 初期版は accuracy {score(legacy.accuracy)} / macro F1 {score(legacy.macro_f1)} / 誤り recall {score(legacy.binary_recall)}。",
        f"- 最新版は accuracy {score(current.accuracy)} / macro F1 {score(current.macro_f1)} / 誤り recall {score(current.binary_recall)}。",
        "- 改善が大きかったことをこのスライドで強く伝える。",
        "",
        "## Slide 8",
        f"- 最新版の主な誤分類は {', '.join(top_current)}。",
        "- ほぼ正確の境界と、誤りを不正確へ落とすケースが残課題。",
        "",
        "## Slide 9",
        "- 精度 84% まで来たことで、説明付き確認支援ツールとしての実用性が高まったと話す。",
        "",
        "## Slide 10",
        "- shared dataset 集約、5 plots、CSV/XLSX bundle 化など運用面の工夫を説明する。",
        "",
        "## Slide 11",
        "- 次の改善テーマは 誤り -> 不正確 の圧縮と、ほぼ正確の境界調整。",
        "",
        "## Slide 12",
        f"- summary_metrics.png では Accuracy {score(current.accuracy)}、Binary F1 {score(current.binary_f1)} を確認する。",
        "",
        "## Slide 13",
        f"- per_class_metrics.png では ほぼ正確 recall {score(current.per_class['ほぼ正確']['recall'])} が弱点であることを確認する。",
        "",
        "## Slide 14",
        "- confusion_matrix.png では ほぼ正確 -> 正確 6件、誤り -> 不正確 4件に注目する。",
        "",
        "## Slide 15",
        f"- evaluation_overview.png では mismatches {current.mismatch_total}件、sample_count 100、skipped 0 を確認する。",
        "",
        "## Slide 16-19",
        "- 最後は参考資料として図を単独表示し、質疑応答で使えるようにする。",
        "",
    ]
    return "\n".join(lines)


def build_presentation(dataset: DatasetSummary, legacy: EvalSnapshot, current: EvalSnapshot) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    current_report = load_json(CURRENT_VER2_EVAL_PATH)
    current_meta = load_json(CURRENT_VER2_PREDICTIONS_PATH)["meta"]
    weighted_f1 = float(current_report["multiclass"]["weighted_avg"]["f1"])
    macro_precision = float(current_report["multiclass"]["macro_avg"]["precision"])
    macro_recall = float(current_report["multiclass"]["macro_avg"]["recall"])
    case_count = int(current_meta.get("case_count") or dataset.case_count)

    current_examples = representative_mismatch_examples(CURRENT_VER2_CSV_PATH, CURRENT_VER2_EVAL_PATH)
    current_flow = [f"{truth} -> {pred}: {count}件" for (truth, pred), count in current.mismatch_counts[:4]]

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    hero.fill.solid()
    hero.fill.fore_color.rgb = NAVY
    hero.line.fill.background()
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(0.98), Inches(0.22), Inches(5.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background()
    add_textbox(slide, Inches(1.15), Inches(1.2), Inches(5.0), Inches(0.35), "datamix", 22, WHITE, True)
    add_textbox(slide, Inches(1.15), Inches(1.62), Inches(6.9), Inches(0.8), "生成AI・機械学習エンジニア育成講座\n卒業発表", 22, WHITE, True)
    add_textbox(slide, Inches(1.15), Inches(3.0), Inches(6.9), Inches(0.6), "フェイクニュースチェッカー", 28, WHITE, True)
    add_textbox(slide, Inches(1.15), Inches(4.82), Inches(6.2), Inches(0.28), "20260401 本編資料 統合版", 15, WHITE, False)
    add_textbox(slide, Inches(1.15), Inches(5.2), Inches(6.2), Inches(0.35), "Fake News Checker / 2026年4月11日時点", 16, WHITE, False)
    add_cover_footer(slide)
    side = add_panel(slide, Inches(9.48), Inches(1.2), Inches(2.95), Inches(4.9), fill_rgb=SAND, line_rgb=SAND)
    side.line.fill.background()
    add_textbox(slide, Inches(9.82), Inches(1.62), Inches(2.2), Inches(0.28), "今回の更新点", 15, SUBTLE, True)
    add_bullets(
        slide,
        Inches(9.72),
        Inches(2.0),
        Inches(2.35),
        Inches(3.45),
        [
            "本編の章立てを維持",
            "最新 Ver2 rerun を反映",
            "評価画像を差し替え",
            "改善点を再整理",
            "質疑向け参考図も更新",
        ],
        font_size=15,
        color=INK,
    )

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "1. 課題と選んだ理由", "誤った情報に踊らされず、正確な情報へ近づく支援を作る")
    add_panel(slide, Inches(0.72), Inches(1.52), Inches(5.85), Inches(4.95), fill_rgb=WHITE)
    add_panel(slide, Inches(6.76), Inches(1.52), Inches(5.85), Inches(4.95), fill_rgb=MINT)
    add_textbox(slide, Inches(0.98), Inches(1.82), Inches(2.3), Inches(0.3), "社会的な背景", 20, TEAL, True)
    add_bullets(
        slide,
        Inches(0.95),
        Inches(2.18),
        Inches(5.1),
        Inches(3.6),
        [
            "オールドメディアへの不信感が高まる一方で、X や YouTube など SNS の台頭により、多様な情報が大量に流通している",
            "兵庫県知事の騒動のように、報道の偏りや受け手の分断が可視化された事例が続き、確認支援の必要性を強く感じた",
            "戦時中にマスメディアが世論形成に大きく関与した歴史も踏まえ、情報をうのみにしない仕組みが必要だと考えた",
        ],
        font_size=17,
    )
    add_textbox(slide, Inches(7.03), Inches(1.82), Inches(2.6), Inches(0.3), "個人的に選んだ理由", 20, TEAL, True)
    add_bullets(
        slide,
        Inches(7.0),
        Inches(2.18),
        Inches(5.1),
        Inches(3.55),
        [
            "昔、誤った情報に騙されて痛い出費をした経験があり、真偽を自分で確かめる重要性を痛感した",
            "強い言葉や断定調に流されず、自分の頭で確認する補助があれば被害を減らせると考えた",
            "だからこそ、断定の代行ではなく、正確な情報へ近づくための確認支援ツールを作りたいと考えた",
        ],
        font_size=17,
    )
    add_metric_card(slide, Inches(6.98), Inches(5.45), Inches(1.55), Inches(0.82), "入力", "URL / 本文", TEAL)
    add_metric_card(slide, Inches(8.76), Inches(5.45), Inches(1.45), Inches(0.82), "出力", "5区分", ORANGE)
    add_metric_card(slide, Inches(10.42), Inches(5.45), Inches(1.88), Inches(0.82), "目標", "確認支援", GREEN)
    add_footer(slide, FOOTER_NOTE)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "2. 使用したデータについて", "balanced 100件の real_article_dataset_v2 を使って最新 Ver2 を評価")
    add_panel(slide, Inches(0.72), Inches(1.62), Inches(3.5), Inches(3.0), fill_rgb=WHITE)
    add_panel(slide, Inches(4.45), Inches(1.62), Inches(3.15), Inches(3.0), fill_rgb=MINT)
    add_panel(slide, Inches(7.83), Inches(1.62), Inches(4.76), Inches(3.0), fill_rgb=WHITE)
    add_textbox(slide, Inches(0.98), Inches(1.92), Inches(2.3), Inches(0.3), "dataset の概要", 18, TEAL, True)
    add_bullets(
        slide,
        Inches(0.95),
        Inches(2.24),
        Inches(2.95),
        Inches(1.95),
        [
            "ファイル: real_article_dataset_v2.json",
            f"件数: {dataset.case_count}",
            f"curated_on: {dataset.curated_on}",
            "5区分は各20件で balanced",
        ],
        font_size=16,
    )
    add_textbox(slide, Inches(4.72), Inches(1.92), Inches(2.0), Inches(0.3), "スキーマ", 18, TEAL, True)
    add_bullets(
        slide,
        Inches(4.68),
        Inches(2.24),
        Inches(2.4),
        Inches(1.95),
        [f"`{field}`" for field in dataset.fields],
        font_size=16,
    )
    add_textbox(slide, Inches(8.1), Inches(1.92), Inches(2.0), Inches(0.3), "ドメイン分布", 18, TEAL, True)
    add_bullets(
        slide,
        Inches(8.06),
        Inches(2.24),
        Inches(4.0),
        Inches(1.95),
        [
            f"一般: {dataset.domain_counts['一般']}",
            f"医療: {dataset.domain_counts['医療']}",
            f"金融: {dataset.domain_counts['金融']}",
            f"災害: {dataset.domain_counts['災害']}",
            f"政治: {dataset.domain_counts['政治']}",
        ],
        font_size=16,
    )
    add_panel(slide, Inches(0.76), Inches(4.95), Inches(11.82), Inches(1.15), fill_rgb=SAND)
    add_bullets(
        slide,
        Inches(0.98),
        Inches(5.18),
        Inches(11.3),
        Inches(0.6),
        [
            "評価対象は記事全文ではなく analysis_text に要約した中心命題で、expected_verdict と predicted_verdict を直接比較できるようにしている",
        ],
        font_size=17,
    )
    add_footer(slide, FOOTER_NOTE)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "3. 初期仮説・分析アプローチ", "ローカル一次判定と Gemini の外部根拠比較を組み合わせれば、境界ケースを改善できると考えた")
    add_panel(slide, Inches(0.74), Inches(1.58), Inches(11.84), Inches(1.82), fill_rgb=SAND)
    add_bullets(
        slide,
        Inches(0.98),
        Inches(1.9),
        Inches(11.25),
        Inches(1.1),
        [
            "仮説1: 出典の有無、著者、日付、リンク数、表現の強さなどのヒューリスティックで危険度の目安は作れる",
            "仮説2: 境界ケースは Gemini の google_search と url_context を使った外部根拠比較を加えることで改善できる",
            "仮説3: 書き振り評価は真偽判定と分離した方が、煽り表現だけで誤判定するリスクを下げられる",
        ],
        font_size=17,
    )
    flow_y = Inches(4.15)
    flow_w = Inches(2.0)
    add_process_box(slide, Inches(0.78), flow_y, flow_w, Inches(1.38), "入力", "URL または本文", TEAL)
    add_connector(slide, Inches(2.9), flow_y + Inches(0.46), Inches(0.35), Inches(0.34), LINE)
    add_process_box(slide, Inches(3.18), flow_y, flow_w, Inches(1.38), "ローカル一次判定", "risk / confidence / domain を推定", ORANGE)
    add_connector(slide, Inches(5.3), flow_y + Inches(0.46), Inches(0.35), Inches(0.34), LINE)
    add_process_box(slide, Inches(5.58), flow_y, flow_w, Inches(1.38), "Gemini 比較", "外部根拠と反証を照合", GREEN)
    add_connector(slide, Inches(7.7), flow_y + Inches(0.46), Inches(0.35), Inches(0.34), LINE)
    add_process_box(slide, Inches(7.98), flow_y, flow_w, Inches(1.38), "書き振り評価", "煽りや断定を別軸で評価", ORANGE)
    add_connector(slide, Inches(10.1), flow_y + Inches(0.46), Inches(0.35), Inches(0.34), LINE)
    add_process_box(slide, Inches(10.38), flow_y, Inches(2.0), Inches(1.38), "出力", "5区分・理由・確認リンク", RED)
    add_footer(slide, FOOTER_NOTE)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "3. プログラムの実行例", "入力から判定結果の表示までを、実際の画面で確認する")
    add_screenshot_card(slide, SCREENSHOT_INPUT, Inches(0.58), Inches(1.62), Inches(4.05), Inches(4.6), "1. 入力画面", TEAL)
    add_screenshot_card(slide, SCREENSHOT_LOADING, Inches(4.64), Inches(1.62), Inches(4.05), Inches(4.6), "2. 判定中", ORANGE)
    add_screenshot_card(slide, SCREENSHOT_RESULT, Inches(8.7), Inches(1.62), Inches(4.05), Inches(4.6), "3. 結果表示", GREEN)
    add_panel(slide, Inches(0.82), Inches(6.36), Inches(11.75), Inches(0.44), fill_rgb=SAND, line_rgb=SAND)
    add_textbox(
        slide,
        Inches(1.0),
        Inches(6.46),
        Inches(11.2),
        Inches(0.18),
        "URL または本文を入力すると、判定中表示を経て、要注意度・5区分の判定・理由が返る。",
        15,
        INK,
        False,
        PP_ALIGN.CENTER,
    )
    add_footer(slide, FOOTER_NOTE)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "3. 出力の見方", "判定ラベルだけでなく、理由と確認補助情報も返す設計")
    add_panel(slide, Inches(0.68), Inches(1.56), Inches(7.1), Inches(5.34), fill_rgb=WHITE)
    add_panel(slide, Inches(8.0), Inches(1.56), Inches(4.62), Inches(2.82), fill_rgb=WHITE)
    add_panel(slide, Inches(8.0), Inches(4.54), Inches(4.62), Inches(2.36), fill_rgb=MINT)
    add_picture_contained(slide, SCREENSHOT_RESULT, Inches(0.88), Inches(1.82), Inches(6.72), Inches(4.86))
    add_picture_contained(slide, SCREENSHOT_REASON, Inches(8.18), Inches(1.78), Inches(4.24), Inches(2.38))
    add_textbox(slide, Inches(8.28), Inches(4.82), Inches(2.2), Inches(0.28), "この画面で伝えたいこと", 18, TEAL, True)
    add_bullets(
        slide,
        Inches(8.24),
        Inches(5.18),
        Inches(3.9),
        Inches(1.35),
        [
            "要注意度は 0〜100% で直感的に表示する",
            "判定は 5 区分で返し、白黒だけに寄せない",
            "理由は短い箇条書きで、利用者が次の確認行動を取りやすくする",
        ],
        font_size=15,
    )
    add_footer(slide, FOOTER_NOTE)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "4. 分析結果", "2026-04-01 の初期版から、2026-04-11 の最新 rerun で大幅に改善した")
    add_picture_contained(slide, plot_path(current.plots_dir, "evaluation_dashboard.png"), Inches(0.72), Inches(1.62), Inches(7.2), Inches(5.1))
    add_panel(slide, Inches(8.1), Inches(1.62), Inches(4.25), Inches(5.05), fill_rgb=SAND)
    add_metric_card(slide, Inches(8.34), Inches(1.96), Inches(1.65), Inches(0.82), "Accuracy", f"{score(legacy.accuracy)} -> {score(current.accuracy)}", TEAL)
    add_metric_card(slide, Inches(10.1), Inches(1.96), Inches(1.65), Inches(0.82), "Macro F1", f"{score(legacy.macro_f1)} -> {score(current.macro_f1)}", ORANGE)
    add_metric_card(slide, Inches(8.34), Inches(3.0), Inches(1.65), Inches(0.82), "誤り Recall", f"{score(legacy.binary_recall)} -> {score(current.binary_recall)}", RED)
    add_metric_card(slide, Inches(10.1), Inches(3.0), Inches(1.65), Inches(0.82), "Binary F1", score(current.binary_f1), GREEN)
    add_metric_card(slide, Inches(8.34), Inches(4.04), Inches(1.65), Inches(0.82), "Mismatches", f"{legacy.mismatch_total} -> {current.mismatch_total}", TEAL)
    add_metric_card(slide, Inches(10.1), Inches(4.04), Inches(1.65), Inches(0.82), "All cases", f"{case_count}/{case_count}", ORANGE)
    add_bullets(
        slide,
        Inches(8.28),
        Inches(5.1),
        Inches(3.65),
        Inches(1.15),
        [
            "Accuracy は 0.370 から 0.840 まで改善",
            "誤り recall も 0.000 から 0.800 へ回復",
            "残る課題は『ほぼ正確』境界と誤りの押し切り",
        ],
        font_size=15,
    )
    add_footer(slide, FOOTER_NOTE)

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "4. 考察", "現在の Ver2 はかなり実用的だが、境界ラベルの調整はまだ必要")
    add_picture_contained(slide, plot_path(current.plots_dir, "confusion_matrix.png"), Inches(0.75), Inches(1.74), Inches(5.8), Inches(4.95))
    add_panel(slide, Inches(6.72), Inches(1.74), Inches(5.55), Inches(4.95), fill_rgb=WHITE)
    add_textbox(slide, Inches(7.0), Inches(2.0), Inches(2.8), Inches(0.3), "主な誤判定", 20, ORANGE, True)
    add_bullets(slide, Inches(6.96), Inches(2.34), Inches(4.85), Inches(1.25), current_flow, font_size=17)
    add_textbox(slide, Inches(7.0), Inches(4.0), Inches(2.8), Inches(0.3), "代表ケース", 20, TEAL, True)
    add_bullets(slide, Inches(6.96), Inches(4.34), Inches(4.85), Inches(1.3), current_examples, font_size=15)
    add_footer(slide, FOOTER_NOTE)

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "5. ビジネスインパクト", "最新評価では、説明付き確認支援ツールとしてかなり使える水準に近づいた")
    add_metric_card(slide, Inches(0.88), Inches(1.82), Inches(3.55), Inches(1.02), "個人向け", "投稿前セルフチェック", TEAL)
    add_metric_card(slide, Inches(4.88), Inches(1.82), Inches(3.55), Inches(1.02), "組織向け", "一次仕分けと優先順位付け", ORANGE)
    add_metric_card(slide, Inches(8.88), Inches(1.82), Inches(3.48), Inches(1.02), "教育向け", "メディアリテラシー支援", GREEN)
    add_panel(slide, Inches(0.88), Inches(3.2), Inches(11.45), Inches(3.0), fill_rgb=SAND)
    add_bullets(
        slide,
        Inches(1.15),
        Inches(3.5),
        Inches(10.9),
        Inches(2.2),
        [
            f"最新 rerun では Accuracy {pct(current.accuracy)}、Binary F1 {score(current.binary_f1)} まで到達し、『全く使えない』段階は抜けた",
            "5区分 verdict に加え、理由・確認リンク・attention score を返すので、利用者が次の確認行動を取りやすい",
            "自治体、学校、メディア運用では『まず人が見るべき候補』を絞る補助として使える",
            "ただし最終断定を完全自動化するには、ほぼ正確境界と誤り recall をもう一段伸ばす必要がある",
        ],
        font_size=18,
    )
    add_footer(slide, FOOTER_NOTE)

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "6. 工夫点", "評価可能な形に落とし込み、レビューしやすい運用まで整えた")
    add_panel(slide, Inches(0.74), Inches(1.62), Inches(5.9), Inches(4.95), fill_rgb=WHITE)
    add_panel(slide, Inches(6.86), Inches(1.62), Inches(5.7), Inches(4.95), fill_rgb=MINT)
    add_textbox(slide, Inches(1.0), Inches(1.92), Inches(2.2), Inches(0.3), "設計面の工夫", 20, TEAL, True)
    add_bullets(
        slide,
        Inches(0.98),
        Inches(2.26),
        Inches(5.1),
        Inches(2.15),
        [
            "真偽判定と書き振り評価を分離し、煽り表現だけで真偽を決めない設計にした",
            "dataset_runner / evaluation / plot をつなぎ、同じ dataset で改善前後を比較できるようにした",
            "shared testdata を repo root に集約し、Ver2 / Ver3 / Ver4 で同じデータを参照できるようにした",
        ],
        font_size=17,
    )
    add_textbox(slide, Inches(1.0), Inches(4.78), Inches(2.8), Inches(0.3), "運用面の工夫", 20, ORANGE, True)
    add_bullets(
        slide,
        Inches(0.98),
        Inches(5.1),
        Inches(5.05),
        Inches(1.0),
        [
            "evaluation_outputs では 5 plots を必須で出す運用にした",
            "predictions JSON / eval JSON / plots / CSV / XLSX を同じ timestamp フォルダへまとめた",
            "プレゼン資料もスクリプトで再生成し、最新 rerun をすぐ反映できるようにした",
        ],
        font_size=17,
    )
    positions = [
        (Inches(7.08), Inches(1.95), Inches(2.2), Inches(1.55), "evaluation_dashboard.png"),
        (Inches(9.46), Inches(1.95), Inches(2.2), Inches(1.55), "confusion_matrix.png"),
        (Inches(7.08), Inches(3.72), Inches(2.2), Inches(1.55), "summary_metrics.png"),
        (Inches(9.46), Inches(3.72), Inches(2.2), Inches(1.55), "per_class_metrics.png"),
    ]
    for left, top, width, height, filename in positions:
        add_labeled_image(slide, plot_path(current.plots_dir, filename), left, top, width, height, filename)
    add_textbox(slide, Inches(7.1), Inches(5.5), Inches(4.2), Inches(0.26), "evaluation_overview.png も同じ bundle に保存", 12, SUBTLE, True)
    add_footer(slide, FOOTER_NOTE)

    # Slide 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "6. 今後さらに実施してみたいこと", "最新評価を踏まえると、次の焦点はかなり明確になった")
    add_panel(slide, Inches(0.76), Inches(1.72), Inches(5.82), Inches(4.92), fill_rgb=WHITE)
    add_panel(slide, Inches(6.82), Inches(1.72), Inches(5.82), Inches(4.92), fill_rgb=MINT)
    add_textbox(slide, Inches(1.02), Inches(2.02), Inches(2.8), Inches(0.3), "モデル改善", 20, TEAL, True)
    add_bullets(
        slide,
        Inches(1.0),
        Inches(2.34),
        Inches(5.05),
        Inches(3.4),
        [
            "誤り -> 不正確 の 4 件を優先して潰し、誤り recall をもう一段引き上げる",
            "ほぼ正確 -> 正確 6 件、ほぼ正確 -> 不正確 2 件を材料に、境界ルールを整理する",
            "Ver4 を主軸に改善しつつ、安全な修正は Ver2 / Ver3 に追従させる",
        ],
        font_size=17,
    )
    add_textbox(slide, Inches(7.08), Inches(2.02), Inches(3.6), Inches(0.3), "データと運用の拡張", 20, ORANGE, True)
    add_bullets(
        slide,
        Inches(7.05),
        Inches(2.34),
        Inches(5.05),
        Inches(3.4),
        [
            "CSV / XLSX を使って誤判定をケース単位で見直し、ルール修正の根拠を増やす",
            "時事性の高いケースや SNS 投稿に近いデータを追加し、実運用に近い評価へ広げる",
            "発表では『今できていること』と『まだできていないこと』を分けて伝える",
        ],
        font_size=17,
    )
    add_footer(slide, FOOTER_NOTE)

    # Slide 12
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "図の説明: summary_metrics.png", "全体指標はかなり改善したが、誤り recall にはまだ伸びしろがある")
    add_picture_contained(slide, plot_path(current.plots_dir, "summary_metrics.png"), Inches(0.78), Inches(1.68), Inches(7.0), Inches(5.0))
    add_panel(slide, Inches(8.02), Inches(1.72), Inches(4.28), Inches(4.9), fill_rgb=SAND)
    add_textbox(slide, Inches(8.28), Inches(2.0), Inches(2.8), Inches(0.3), "読み取りポイント", 19, TEAL, True)
    add_bullets(
        slide,
        Inches(8.22),
        Inches(2.34),
        Inches(3.6),
        Inches(2.8),
        [
            f"Accuracy は {score(current.accuracy)}、Macro F1 は {score(current.macro_f1)} まで改善",
            f"Weighted F1 も {score(weighted_f1)} で、全体として安定してきた",
            f"Binary Precision {score(current.binary_precision)} に対し、Binary Recall は {score(current.binary_recall)}",
            "つまり『誤り』を出したときの精度は高いが、まだ 20% ほど取りこぼしがある",
        ],
        font_size=17,
    )
    add_footer(slide, FOOTER_NOTE)

    # Slide 13
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "図の説明: per_class_metrics.png", "クラスごとに見ると、保留と不正確は強く、ほぼ正確の回収が弱い")
    add_picture_contained(slide, plot_path(current.plots_dir, "per_class_metrics.png"), Inches(0.72), Inches(1.68), Inches(7.45), Inches(5.02))
    add_panel(slide, Inches(8.35), Inches(1.72), Inches(3.95), Inches(4.9), fill_rgb=WHITE)
    add_textbox(slide, Inches(8.62), Inches(2.0), Inches(2.8), Inches(0.3), "読み取りポイント", 19, TEAL, True)
    add_bullets(
        slide,
        Inches(8.58),
        Inches(2.34),
        Inches(3.2),
        Inches(2.85),
        [
            f"判断保留は precision / recall / F1 がすべて {score(current.per_class['判断保留']['f1'])}",
            f"不正確は recall {score(current.per_class['不正確']['recall'])} で安定している",
            f"ほぼ正確は recall {score(current.per_class['ほぼ正確']['recall'])} で、ここが現在の弱点",
            f"誤りは recall {score(current.per_class['誤り']['recall'])} まで改善したが、まだ上げたい",
        ],
        font_size=17,
    )
    add_footer(slide, FOOTER_NOTE)

    # Slide 14
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "図の説明: confusion_matrix.png", "どのラベルへ流れているかを見ると、今の残課題がはっきり見える")
    add_picture_contained(slide, plot_path(current.plots_dir, "confusion_matrix.png"), Inches(0.72), Inches(1.68), Inches(6.8), Inches(5.0))
    add_panel(slide, Inches(7.8), Inches(1.72), Inches(4.5), Inches(4.9), fill_rgb=MINT)
    add_textbox(slide, Inches(8.05), Inches(2.0), Inches(3.1), Inches(0.3), "読み取りポイント", 19, TEAL, True)
    add_bullets(
        slide,
        Inches(8.0),
        Inches(2.34),
        Inches(3.75),
        Inches(2.85),
        [
            "正確 20 件のうち 16 件は正解、4 件は ほぼ正確 に落ちた",
            "ほぼ正確 20 件のうち 6 件が 正確、2 件が 不正確 に流れた",
            "誤り 20 件のうち 16 件は正解、4 件は 不正確 に落ちた",
            "改善の中心は『ほぼ正確』境界と『誤り』の押し切りであることが分かる",
        ],
        font_size=17,
    )
    add_footer(slide, FOOTER_NOTE)

    # Slide 15
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "図の説明: evaluation_overview.png", "集計サマリーからも、かなり改善したことと残課題が同時に確認できる")
    add_picture_contained(slide, plot_path(current.plots_dir, "evaluation_overview.png"), Inches(1.2), Inches(1.65), Inches(4.8), Inches(4.95))
    add_panel(slide, Inches(6.35), Inches(1.72), Inches(5.95), Inches(4.9), fill_rgb=SAND)
    add_textbox(slide, Inches(6.65), Inches(2.0), Inches(3.0), Inches(0.3), "読み取りポイント", 19, TEAL, True)
    add_bullets(
        slide,
        Inches(6.6),
        Inches(2.34),
        Inches(5.15),
        Inches(2.75),
        [
            f"sample_count は {case_count}、全件処理できており skipped は 0",
            f"mismatches は {current.mismatch_total} 件で、初期版 63 件から大幅に減少した",
            f"Macro Avg precision {score(macro_precision)} / recall {score(macro_recall)} / f1 {score(current.macro_f1)}",
            "総合性能はかなり伸びたが、mismatch 16 件をどう詰めるかが次のテーマになる",
        ],
        font_size=17,
    )
    add_footer(slide, FOOTER_NOTE)

    # Slide 16
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "参考資料: confusion_matrix.png", "Ver2 / evaluation_outputs/20260411-0509/plots")
    add_picture_contained(slide, plot_path(current.plots_dir, "confusion_matrix.png"), Inches(1.0), Inches(1.5), Inches(11.2), Inches(5.3))
    add_footer(slide, FOOTER_NOTE)

    # Slide 17
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "参考資料: evaluation_overview.png", "Ver2 / evaluation_outputs/20260411-0509/plots")
    add_picture_contained(slide, plot_path(current.plots_dir, "evaluation_overview.png"), Inches(3.15), Inches(1.45), Inches(6.9), Inches(5.45))
    add_footer(slide, FOOTER_NOTE)

    # Slide 18
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "参考資料: per_class_metrics.png", "Ver2 / evaluation_outputs/20260411-0509/plots")
    add_picture_contained(slide, plot_path(current.plots_dir, "per_class_metrics.png"), Inches(0.76), Inches(1.54), Inches(11.85), Inches(5.35))
    add_footer(slide, FOOTER_NOTE)

    # Slide 19
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "参考資料: summary_metrics.png", "Ver2 / evaluation_outputs/20260411-0509/plots")
    add_picture_contained(slide, plot_path(current.plots_dir, "summary_metrics.png"), Inches(0.85), Inches(1.58), Inches(11.65), Inches(5.3))
    add_footer(slide, FOOTER_NOTE)

    return prs


def main() -> None:
    dataset = build_dataset_summary(DATASET_PATH)
    legacy = build_eval_snapshot("Ver2 tune4", LEGACY_VER2_EVAL_PATH, LEGACY_VER2_PREDICTIONS_PATH, LEGACY_VER2_PLOTS_DIR)
    current = build_eval_snapshot("Ver2 latest", CURRENT_VER2_EVAL_PATH, CURRENT_VER2_PREDICTIONS_PATH, CURRENT_VER2_PLOTS_DIR)
    prs = build_presentation(dataset, legacy, current)
    prs.save(OUTPUT_PPTX)
    OUTPUT_NOTES.write_text(build_notes(dataset, legacy, current), encoding="utf-8")
    print(OUTPUT_PPTX)
    print(OUTPUT_NOTES)


if __name__ == "__main__":
    main()
