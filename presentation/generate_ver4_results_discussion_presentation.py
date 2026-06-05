from __future__ import annotations

import json
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PRESENTATION_DIR = Path(__file__).resolve().parent
PROJECT_ROOT = PRESENTATION_DIR.parent

DATASET_PATH = PROJECT_ROOT / "testdata" / "shared" / "real_article_dataset_v2.json"
EVAL_PATH = PROJECT_ROOT / "Ver4" / "evaluation_outputs" / "20260427-0304" / "eval_real_article_dataset_v2_use_gpt_gemini.json"
PLOTS_DIR = PROJECT_ROOT / "Ver4" / "evaluation_outputs" / "20260427-0304" / "plots"

OUTPUT_PPTX = PRESENTATION_DIR / "ver4_results_discussion_only_20260429.pptx"
OUTPUT_NOTES = PRESENTATION_DIR / "ver4_results_discussion_only_20260429_notes.md"

FONT = "Yu Gothic"
BG = RGBColor(247, 248, 246)
INK = RGBColor(27, 35, 43)
SUBTLE = RGBColor(88, 99, 110)
WHITE = RGBColor(255, 255, 255)
NAVY = RGBColor(24, 48, 72)
TEAL = RGBColor(15, 124, 128)
GREEN = RGBColor(61, 143, 99)
GOLD = RGBColor(191, 151, 48)
ORANGE = RGBColor(207, 115, 47)
RED = RGBColor(182, 63, 63)
LINE = RGBColor(214, 220, 224)
MINT = RGBColor(230, 243, 240)
BLUE_SOFT = RGBColor(232, 239, 248)
GREEN_SOFT = RGBColor(232, 244, 235)
GOLD_SOFT = RGBColor(251, 246, 226)
ORANGE_SOFT = RGBColor(250, 238, 227)
RED_SOFT = RGBColor(249, 232, 232)


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def pct(value: float) -> str:
    return f"{value * 100:.1f}%"


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_text(slide, left, top, width, height, text="", size=20, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(3)
    tf.margin_right = Pt(3)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if text:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        p.alignment = align
    return box


def add_header(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, Inches(0.62), Inches(0.34), Inches(10.8), Inches(0.42), title, 25, NAVY, True)
    if subtitle:
        add_text(slide, Inches(0.64), Inches(0.81), Inches(11.1), Inches(0.26), subtitle, 10, SUBTLE)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(1.1), Inches(12.05), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def add_footer(slide, page: int) -> None:
    add_text(slide, Inches(0.65), Inches(7.05), Inches(7.2), Inches(0.18), "Fake News Checker Ver4 | 4. 分析結果・考察", 8, SUBTLE)
    add_text(slide, Inches(12.25), Inches(7.05), Inches(0.55), Inches(0.18), str(page), 8, SUBTLE, False, PP_ALIGN.RIGHT)


def add_panel(slide, left, top, width, height, fill=WHITE, line=LINE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_metric(slide, left, top, width, height, label: str, value: str, accent=TEAL, note: str | None = None):
    add_panel(slide, left, top, width, height)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.08))
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()
    add_text(slide, left + Inches(0.12), top + Inches(0.2), width - Inches(0.24), Inches(0.2), label, 11, SUBTLE, True, PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.12), top + Inches(0.47), width - Inches(0.24), Inches(0.42), value, 24, accent, True, PP_ALIGN.CENTER)
    if note:
        add_text(slide, left + Inches(0.12), top + Inches(0.93), width - Inches(0.24), Inches(0.23), note, 9, SUBTLE, False, PP_ALIGN.CENTER)


def add_bullets(slide, left, top, width, height, items: list[str], size=16, color=INK):
    box = add_text(slide, left, top, width, height)
    tf = box.text_frame
    tf.clear()
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(5)
        p.line_spacing = 1.1


def add_picture_fit(slide, image_path: Path, left, top, width, height):
    add_panel(slide, left, top, width, height, WHITE)
    pic = slide.shapes.add_picture(str(image_path), left + Inches(0.08), top + Inches(0.08))
    max_w = width - Inches(0.16)
    max_h = height - Inches(0.16)
    scale = min(max_w / pic.width, max_h / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = left + int((width - pic.width) / 2)
    pic.top = top + int((height - pic.height) / 2)


def add_result_row(slide, top, label: str, precision: float, recall: float, f1: float, accent):
    add_panel(slide, Inches(0.92), top, Inches(11.45), Inches(0.5), WHITE)
    add_text(slide, Inches(1.18), top + Inches(0.13), Inches(1.4), Inches(0.16), label, 12, accent, True)
    add_text(slide, Inches(3.1), top + Inches(0.13), Inches(1.4), Inches(0.16), f"Precision {pct(precision)}", 10, SUBTLE)
    add_text(slide, Inches(5.25), top + Inches(0.13), Inches(1.4), Inches(0.16), f"Recall {pct(recall)}", 10, SUBTLE)
    add_text(slide, Inches(7.15), top + Inches(0.13), Inches(1.4), Inches(0.16), f"F1 {pct(f1)}", 10, SUBTLE)


def build_notes(eval_report: dict, mismatch_patterns: list[tuple[tuple[str, str], int]]) -> str:
    top_patterns = [f"{truth} -> {pred}: {count}件" for (truth, pred), count in mismatch_patterns]
    return "\n".join(
        [
            "# Ver4 Results Discussion Notes",
            "",
            "- generated_on: 2026-04-29",
            f"- eval: {EVAL_PATH.relative_to(PROJECT_ROOT)}",
            "",
            "## Slide 1",
            "- 分析結果・考察のみの抜粋資料。",
            "",
            "## Slide 2",
            f"- Ver4の最新評価はAccuracy {pct(eval_report['multiclass']['accuracy'])}、Macro F1 {pct(eval_report['multiclass']['macro_avg']['f1'])}。",
            f"- ミスマッチは{len(eval_report['mismatches'])}件。",
            "",
            "## Slide 3",
            "- ラベル別では、正確・判断保留・誤りが特に安定している。",
            "- ほぼ正確と不正確は境界判定が難しく、改善余地がある。",
            "",
            "## Slide 4",
            "- 混同行列から、主な誤分類はほぼ正確と正確、不正確とほぼ正確の境界に集中している。",
            *[f"- {pattern}" for pattern in top_patterns],
            "",
            "## Slide 5",
            "- 考察として、Ver4は明確な正誤や判断保留には強い。",
            "- 一方、中間ラベルは人間でも揺れやすく、構造的な境界ルールが必要。",
            "",
            "## Slide 6",
            "- 今後はblind testで汎化性能を確認し、固有名詞ごとの個別対応ではなく境界判定ルールを改善する。",
        ]
    ) + "\n"


def build_presentation() -> None:
    eval_report = load_json(EVAL_PATH)
    dataset = load_json(DATASET_PATH)
    multiclass = eval_report["multiclass"]
    per_class = multiclass["per_class"]
    mismatch_patterns = Counter((item["truth"], item["pred"]) for item in eval_report["mismatches"]).most_common()
    verdict_counts = Counter(case["expected_verdict"] for case in dataset["cases"])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    add_text(slide, Inches(0.82), Inches(1.28), Inches(10), Inches(0.42), "Fake News Checker Ver4", 25, RGBColor(205, 229, 227), True)
    add_text(slide, Inches(0.82), Inches(1.88), Inches(10.9), Inches(1.05), "4. 分析結果・考察", 46, WHITE, True)
    add_text(slide, Inches(0.86), Inches(3.08), Inches(8.8), Inches(0.55), "最新100件評価から見えた性能と残課題", 20, RGBColor(232, 238, 244))
    add_metric(slide, Inches(0.9), Inches(4.55), Inches(2.45), Inches(1.15), "Accuracy", pct(multiclass["accuracy"]), TEAL)
    add_metric(slide, Inches(3.55), Inches(4.55), Inches(2.45), Inches(1.15), "Macro F1", pct(multiclass["macro_avg"]["f1"]), GREEN)
    add_metric(slide, Inches(6.2), Inches(4.55), Inches(2.45), Inches(1.15), "Mismatches", str(len(eval_report["mismatches"])), ORANGE)
    add_text(slide, Inches(0.9), Inches(6.38), Inches(4.3), Inches(0.22), "2026-04-29", 12, RGBColor(214, 224, 232))

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "評価設計と全体結果", "100件のbalanced datasetで5段階分類を評価")
    add_metric(slide, Inches(0.78), Inches(1.45), Inches(2.15), Inches(1.18), "Sample", str(multiclass.get("sample_count") or eval_report["meta"]["sample_count"]), NAVY, "5ラベル各20件")
    add_metric(slide, Inches(3.18), Inches(1.45), Inches(2.15), Inches(1.18), "Accuracy", pct(multiclass["accuracy"]), TEAL)
    add_metric(slide, Inches(5.58), Inches(1.45), Inches(2.15), Inches(1.18), "Macro F1", pct(multiclass["macro_avg"]["f1"]), GREEN)
    add_metric(slide, Inches(7.98), Inches(1.45), Inches(2.15), Inches(1.18), "Weighted F1", pct(multiclass["weighted_avg"]["f1"]), GOLD)
    add_metric(slide, Inches(10.38), Inches(1.45), Inches(2.15), Inches(1.18), "Skipped", str(eval_report["meta"]["skipped_count"]), RED)
    add_picture_fit(slide, PLOTS_DIR / "summary_metrics.png", Inches(0.78), Inches(3.05), Inches(5.8), Inches(2.95))
    add_panel(slide, Inches(6.92), Inches(3.05), Inches(5.45), Inches(2.95), BLUE_SOFT)
    add_text(slide, Inches(7.2), Inches(3.38), Inches(4.92), Inches(0.3), "結果の読み取り", 20, NAVY, True)
    add_bullets(
        slide,
        Inches(7.18),
        Inches(4.02),
        Inches(4.95),
        Inches(1.35),
        [
            "100件中90件で期待ラベルと一致",
            "5ラベル各20件なので、偏りに頼らない評価",
            "ミスマッチは10件に圧縮",
        ],
        15,
    )
    add_text(slide, Inches(0.9), Inches(6.35), Inches(11.45), Inches(0.28), "全体として、5段階分類でも約90%の性能まで到達した。", 16, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 2)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "ラベル別性能", "安定しているラベルと、境界が難しいラベルを分けて見る")
    add_picture_fit(slide, PLOTS_DIR / "per_class_metrics.png", Inches(0.72), Inches(1.38), Inches(5.75), Inches(4.7))
    colors = {"正確": GREEN, "ほぼ正確": TEAL, "判断保留": GOLD, "不正確": ORANGE, "誤り": RED}
    for i, label in enumerate(["正確", "ほぼ正確", "判断保留", "不正確", "誤り"]):
        row = per_class[label]
        add_result_row(slide, Inches(1.48 + i * 0.58), label, row["precision"], row["recall"], row["f1"], colors[label])
    add_panel(slide, Inches(6.82), Inches(4.76), Inches(5.55), Inches(1.28), MINT)
    add_text(slide, Inches(7.12), Inches(5.05), Inches(5.0), Inches(0.25), "考察", 17, TEAL, True)
    add_text(slide, Inches(7.12), Inches(5.45), Inches(4.95), Inches(0.32), "正確・判断保留・誤りは安定。一方で、ほぼ正確と不正確は境界判定が難しい。", 13, INK)
    add_footer(slide, 3)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "混同行列と主な誤分類", "残ミスマッチ10件は中間ラベルの境界に集中")
    add_picture_fit(slide, PLOTS_DIR / "confusion_matrix.png", Inches(0.72), Inches(1.38), Inches(5.75), Inches(4.7))
    add_text(slide, Inches(6.9), Inches(1.48), Inches(4.8), Inches(0.3), "主なミスマッチ", 21, NAVY, True)
    top = Inches(2.1)
    for (truth, pred), count in mismatch_patterns[:6]:
        add_panel(slide, Inches(6.9), top, Inches(5.15), Inches(0.5), WHITE)
        add_text(slide, Inches(7.16), top + Inches(0.13), Inches(3.2), Inches(0.16), f"{truth} → {pred}", 12, INK, True)
        add_text(slide, Inches(11.08), top + Inches(0.13), Inches(0.62), Inches(0.16), f"{count}件", 12, ORANGE, True, PP_ALIGN.RIGHT)
        top += Inches(0.62)
    add_text(slide, Inches(6.95), Inches(5.9), Inches(5.25), Inches(0.34), "特に「ほぼ正確 → 正確」が4件あり、軽微なズレをどこまで重く見るかが課題。", 13, SUBTLE)
    add_footer(slide, 4)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "考察", "Ver4の強みと残課題")
    strengths = [
        ("明確な正誤に強い", "正確 recall 100.0%、誤り F1 95.0%。根拠と一致・不一致がはっきりするケースは安定。", GREEN, GREEN_SOFT),
        ("判断保留を扱える", "判断保留 F1 97.4%。根拠不足や未確定情報を無理に断定しにくい。", GOLD, GOLD_SOFT),
        ("中間ラベルを表現できる", "二値分類では拾いにくい、軽微差分や一部真実を分けて扱える。", TEAL, MINT),
        ("境界判定は残課題", "正確/ほぼ正確、ほぼ正確/不正確、不正確/誤りの境界はまだ揺れやすい。", ORANGE, ORANGE_SOFT),
    ]
    for i, (title, body, color, fill) in enumerate(strengths):
        left = Inches(0.78 + (i % 2) * 6.05)
        top = Inches(1.55 + (i // 2) * 1.78)
        add_panel(slide, left, top, Inches(5.55), Inches(1.32), fill)
        add_text(slide, left + Inches(0.22), top + Inches(0.22), Inches(5.0), Inches(0.25), title, 17, color, True)
        add_text(slide, left + Inches(0.22), top + Inches(0.64), Inches(5.0), Inches(0.42), body, 12, INK)
    add_panel(slide, Inches(1.08), Inches(5.58), Inches(11.15), Inches(0.82), BLUE_SOFT)
    add_text(slide, Inches(1.35), Inches(5.86), Inches(10.6), Inches(0.25), "結論: Ver4は高精度だが、今後は中間ラベルの境界ルールをさらに磨く必要がある。", 17, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 5)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "今後の改善方針", "既存100件での高性能を、未知データで確認する")
    items = [
        ("blind test 100件を作る", "既存データに合わせすぎていないかを確認する", TEAL),
        ("同じblind setでVer2/3/4/5を比較", "方式の違いを公平に評価する", GREEN),
        ("境界判定ルールを構造化", "固有名詞ごとの個別対応ではなく、一般化できるルールにする", GOLD),
        ("devとblindの両方で確認", "既存性能を落とさず、汎化性能を上げる", ORANGE),
    ]
    for i, (title, body, color) in enumerate(items):
        top = Inches(1.58 + i * 1.05)
        add_panel(slide, Inches(1.02), top, Inches(11.25), Inches(0.72), WHITE)
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(1.23), top + Inches(0.18), Inches(0.34), Inches(0.34))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        add_text(slide, Inches(1.23), top + Inches(0.23), Inches(0.34), Inches(0.16), str(i + 1), 9, WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, Inches(1.78), top + Inches(0.17), Inches(3.2), Inches(0.2), title, 15, color, True)
        add_text(slide, Inches(5.1), top + Inches(0.18), Inches(6.55), Inches(0.2), body, 12, INK)
    add_text(slide, Inches(0.92), Inches(6.3), Inches(11.45), Inches(0.3), "プレゼンでの一言: 既存100件ではVer4が最良。ただし、次はblind testで汎化性能を確認する。", 15, RED, True, PP_ALIGN.CENTER)
    add_footer(slide, 6)

    prs.save(OUTPUT_PPTX)
    OUTPUT_NOTES.write_text(build_notes(eval_report, mismatch_patterns), encoding="utf-8")


if __name__ == "__main__":
    build_presentation()
    print(f"saved: {OUTPUT_PPTX}")
    print(f"saved: {OUTPUT_NOTES}")
