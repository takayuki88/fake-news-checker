from __future__ import annotations

import json
from collections import Counter
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PRESENTATION_DIR = Path(__file__).resolve().parent
PROJECT_ROOT = PRESENTATION_DIR.parent
VER2_ROOT = PROJECT_ROOT / "Ver2"
VER3_ROOT = PROJECT_ROOT / "Ver3"

DATASET_PATH = VER3_ROOT / "testdata" / "real_article_dataset_v2.json"

VER2_EVAL_PATH = VER2_ROOT / "evaluation_outputs" / "20260401-1213" / "eval_real_v2_use_gemini_after_rule_tune4.json"
VER2_PREDICTIONS_PATH = VER2_ROOT / "evaluation_outputs" / "20260401-1213" / "predictions_real_v2_use_gemini_after_rule_tune4.json"
VER2_PLOTS_DIR = VER2_ROOT / "evaluation_outputs" / "20260401-1213" / "plots"

VER3_EVAL_PATH = VER3_ROOT / "evaluation_outputs" / "20260401-1155" / "eval_real_v2_use_gemini_after_rule_tune3.json"
VER3_PREDICTIONS_PATH = VER3_ROOT / "evaluation_outputs" / "20260401-1155" / "predictions_real_v2_use_gemini_after_rule_tune3.json"
VER3_PLOTS_DIR = VER3_ROOT / "evaluation_outputs" / "20260401-1155" / "plots"

EXPORT_BUNDLE_SCRIPT = VER2_ROOT / "scripts" / "export_evaluation_bundle.ps1"

OUTPUT_PPTX = PRESENTATION_DIR / "fake_news_checker_presentation_20260401.pptx"
FALLBACK_OUTPUT_PPTX = PRESENTATION_DIR / "fake_news_checker_presentation_20260401_updated.pptx"
OUTPUT_NOTES = PRESENTATION_DIR / "fake_news_checker_presentation_20260401_notes.md"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(247, 244, 238)
INK = RGBColor(31, 42, 52)
SUBTLE = RGBColor(98, 107, 118)
LINE = RGBColor(217, 211, 202)
WHITE = RGBColor(255, 255, 255)
NAVY = RGBColor(22, 39, 57)
TEAL = RGBColor(19, 118, 122)
ORANGE = RGBColor(201, 109, 49)
GREEN = RGBColor(69, 142, 95)
GOLD = RGBColor(192, 145, 43)
RED = RGBColor(181, 61, 61)
SAND = RGBColor(243, 238, 229)
MINT = RGBColor(234, 243, 241)

FONT_TITLE = "Yu Gothic"
FONT_BODY = "Yu Gothic"

MANDATORY_PLOTS = [
    "evaluation_dashboard.png",
    "confusion_matrix.png",
    "summary_metrics.png",
    "per_class_metrics.png",
    "evaluation_overview.png",
]


@dataclass
class DatasetSummary:
    curated_on: str
    case_count: int
    fields: list[str]
    verdict_counts: Counter
    domain_counts: Counter


@dataclass
class EvalSnapshot:
    label: str
    analysis_datetime: str
    use_gemini: bool
    accuracy: float
    macro_f1: float
    binary_precision: float
    binary_recall: float
    binary_f1: float
    per_class: dict[str, dict]
    predicted_counts: Counter
    mismatch_counts: list[tuple[tuple[str, str], int]]
    mismatch_total: int
    plots_dir: Path


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def build_dataset_summary(path: Path) -> DatasetSummary:
    payload = load_json(path)
    cases = payload["cases"]
    meta = payload["meta"]
    verdict_counts = Counter(case.get("expected_verdict") for case in cases)
    domain_counts = Counter(case.get("expected_domain") for case in cases)
    return DatasetSummary(
        curated_on=str(meta.get("curated_on") or ""),
        case_count=int(meta.get("case_count") or len(cases)),
        fields=list(meta.get("schema_fields") or []),
        verdict_counts=verdict_counts,
        domain_counts=domain_counts,
    )


def build_eval_snapshot(label: str, eval_path: Path, predictions_path: Path, plots_dir: Path) -> EvalSnapshot:
    report = load_json(eval_path)
    predictions = load_json(predictions_path)
    predicted_counts = Counter(record.get("predicted", {}).get("verdict") for record in predictions["records"])
    mismatch_counts = Counter((item["truth"], item["pred"]) for item in report.get("mismatches", []))
    return EvalSnapshot(
        label=label,
        analysis_datetime=str(predictions["meta"].get("analysis_datetime") or ""),
        use_gemini=bool(predictions["meta"].get("use_gemini")),
        accuracy=float(report["multiclass"]["accuracy"]),
        macro_f1=float(report["multiclass"]["macro_avg"]["f1"]),
        binary_precision=float(report["binary_fake_positive"]["precision"]),
        binary_recall=float(report["binary_fake_positive"]["recall"]),
        binary_f1=float(report["binary_fake_positive"]["f1"]),
        per_class=dict(report["multiclass"]["per_class"]),
        predicted_counts=predicted_counts,
        mismatch_counts=mismatch_counts.most_common(5),
        mismatch_total=len(report.get("mismatches", [])),
        plots_dir=plots_dir,
    )


def pct(value: float) -> str:
    return f"{value * 100:.1f}%"


def score(value: float) -> str:
    return f"{value:.3f}"


def set_slide_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_panel(slide, left, top, width, height, fill_rgb=WHITE, line_rgb=LINE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = line_rgb
    shape.line.width = Pt(1.15)
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str = "",
    font_size: int = 20,
    color: RGBColor = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    if text:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = FONT_BODY
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        p.alignment = align
    return box


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    items: list[str],
    font_size: int = 18,
    color: RGBColor = INK,
    spacing: int = 7,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.space_after = Pt(spacing)
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = FONT_BODY
            run.font.size = Pt(font_size)
            run.font.color.rgb = color


def add_header(slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(slide, Inches(0.62), Inches(0.35), Inches(9.2), Inches(0.48), title, 26, INK, True)
    if subtitle:
        add_textbox(slide, Inches(0.64), Inches(0.87), Inches(9.8), Inches(0.32), subtitle, 11, SUBTLE, False)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.63), Inches(1.15), Inches(1.6), Inches(0.07))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def add_footer(slide, note: str = "Fake News Checker | 2026-04-01 時点の運用・評価をもとに作成") -> None:
    add_textbox(slide, Inches(0.65), Inches(7.03), Inches(12.0), Inches(0.2), note, 9, SUBTLE, False, PP_ALIGN.RIGHT)


def add_metric_card(slide, left, top, width, height, label: str, value: str, accent: RGBColor) -> None:
    card = add_panel(slide, left, top, width, height, fill_rgb=WHITE)
    card.line.color.rgb = accent
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.08))
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()
    add_textbox(slide, left + Inches(0.15), top + Inches(0.18), width - Inches(0.3), Inches(0.22), label, 12, SUBTLE, False)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.42), width - Inches(0.3), Inches(0.35), value, 22, INK, True)


def add_process_box(slide, left, top, width, height, title: str, desc: str, fill_rgb: RGBColor) -> None:
    shape = add_panel(slide, left, top, width, height, fill_rgb=fill_rgb, line_rgb=fill_rgb)
    shape.line.color.rgb = fill_rgb
    add_textbox(slide, left + Inches(0.12), top + Inches(0.12), width - Inches(0.24), Inches(0.25), title, 15, WHITE, True)
    desc_box = add_textbox(slide, left + Inches(0.12), top + Inches(0.42), width - Inches(0.24), height - Inches(0.55), desc, 12, WHITE, False)
    desc_box.text_frame.vertical_anchor = MSO_ANCHOR.TOP


def add_connector(slide, left, top, width, height, color: RGBColor) -> None:
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()


def add_labeled_image(slide, image_path: Path, left, top, width, height, label: str) -> None:
    frame = add_panel(slide, left, top, width, height, fill_rgb=WHITE, line_rgb=LINE)
    frame.line.color.rgb = LINE
    slide.shapes.add_picture(str(image_path), left + Inches(0.08), top + Inches(0.08), width=width - Inches(0.16), height=height - Inches(0.42))
    add_textbox(slide, left + Inches(0.1), top + height - Inches(0.27), width - Inches(0.2), Inches(0.15), label, 10, SUBTLE, True, PP_ALIGN.CENTER)


def plot_path(plots_dir: Path, filename: str) -> Path:
    return plots_dir / filename


def build_notes(dataset: DatasetSummary, ver2: EvalSnapshot, ver3: EvalSnapshot) -> str:
    top_ver2 = [f"{truth} -> {pred}: {count}件" for (truth, pred), count in ver2.mismatch_counts[:3]]
    notes = [
        "# Fake News Checker Presentation Notes",
        "",
        f"- generated_on: 2026-04-01",
        f"- dataset: {DATASET_PATH.relative_to(PROJECT_ROOT)}",
        f"- ver2_eval: {VER2_EVAL_PATH.relative_to(PROJECT_ROOT)}",
        f"- ver3_eval: {VER3_EVAL_PATH.relative_to(PROJECT_ROOT)}",
        "",
        "## Slide 1",
        "- 表紙は datamix 卒業発表の見本に合わせたレイアウト。",
        "- タイトルは『フェイクニュースチェッカー』、副題はプレゼン資料であることを明示する。",
        "",
        "## Slide 2",
        "- 課題と選んだ理由では、オールドメディア不信、SNSの台頭、兵庫県知事の騒動、個人的な苦い経験を背景として話す。",
        "- 正確な情報へ近づくための支援ツールを作りたい、という動機にまとめる。",
        "",
        "## Slide 3",
        f"- dataset v2 は {dataset.case_count}件、5区分は各20件で balanced。",
        f"- スキーマは {', '.join(dataset.fields)}。",
        f"- ドメイン分布は 一般{dataset.domain_counts['一般']} / 医療{dataset.domain_counts['医療']} / 金融{dataset.domain_counts['金融']} / 災害{dataset.domain_counts['災害']} / 政治{dataset.domain_counts['政治']}。",
        "",
        "## Slide 4",
        "- 初期仮説は、ローカル一次判定と Gemini の外部根拠比較を組み合わせれば境界ケースを改善できる、というもの。",
        "- 書き振り評価は真偽判定と分離する設計を説明する。",
        "",
        "## Slide 5",
        f"- Ver2 tune4: accuracy {score(ver2.accuracy)}, macro F1 {score(ver2.macro_f1)}, 誤り recall {score(ver2.binary_recall)}。",
        f"- Ver3 tune3: accuracy {score(ver3.accuracy)}, macro F1 {score(ver3.macro_f1)}, 誤り recall {score(ver3.binary_recall)}。",
        "- Ver2 の方が全体指標は少し上だが、誤り recall 0.000 が共通課題であることを話す。",
        "",
        "## Slide 6",
        f"- Ver2 の主な誤判定は {', '.join(top_ver2)}。",
        "- 誤りを『判断保留』や『不正確』へ逃がしている点がボトルネック。",
        "",
        "## Slide 7",
        "- 現状の最適ポジションは自動ジャッジではなく、説明付き確認支援ツール。",
        "- 個人の投稿前チェック、自治体や学校の一次仕分け、教育用途に価値がある。",
        "",
        "## Slide 8",
        "- 工夫点として、真偽判定と書き振り評価の分離、5 plots 運用、CSV/XLSX 出力まで含めた bundle 化を挙げる。",
        f"- まとめ出力は {EXPORT_BUNDLE_SCRIPT.relative_to(PROJECT_ROOT)} を使う。",
        "",
        "## Slide 9",
        "- 今後は Ver2 を中心に、誤り経路の強化、正確ラベル復帰、ケースレビューの継続を行う。",
        "",
        "## Slide 10",
        "- summary_metrics.png を使って、全体指標と binary 指標の弱さを説明する。",
        "",
        "## Slide 11",
        "- per_class_metrics.png を使って、クラスごとの偏りを説明する。",
        "",
        "## Slide 12",
        "- confusion_matrix.png を使って、どのラベルへ誤って流れているかを説明する。",
        "",
        "## Slide 13",
        "- evaluation_overview.png を使って、mismatch 数と集計値を説明する。",
        "",
        "## Slide 14",
        "- 参考資料として confusion_matrix.png を単独スライドで追加。",
        "",
        "## Slide 15",
        "- 参考資料として evaluation_overview.png を単独スライドで追加。",
        "",
        "## Slide 16",
        "- 参考資料として per_class_metrics.png を単独スライドで追加。",
        "",
        "## Slide 17",
        "- 参考資料として summary_metrics.png を単独スライドで追加。",
        "",
    ]
    return "\n".join(notes)


def build_presentation(dataset: DatasetSummary, ver2: EvalSnapshot, ver3: EvalSnapshot) -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    top_ver2 = [f"{truth} -> {pred}: {count}件" for (truth, pred), count in ver2.mismatch_counts[:3]]

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    hero.fill.solid()
    hero.fill.fore_color.rgb = NAVY
    hero.line.fill.background()
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(0.98), Inches(0.22), Inches(5.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background()
    add_textbox(slide, Inches(1.15), Inches(1.2), Inches(5.0), Inches(0.35), "datamix", 22, WHITE, True)
    add_textbox(slide, Inches(1.15), Inches(1.62), Inches(6.7), Inches(0.8), "生成AI・機械学習エンジニア育成講座\n卒業発表", 22, WHITE, True)
    add_textbox(slide, Inches(1.15), Inches(3.05), Inches(6.6), Inches(0.6), "フェイクニュースチェッカー", 28, WHITE, True)
    add_textbox(slide, Inches(1.15), Inches(4.92), Inches(4.2), Inches(0.25), "プレゼン資料", 14, RGBColor(213, 219, 225), False)
    add_textbox(slide, Inches(1.15), Inches(5.28), Inches(5.0), Inches(0.35), "Fake News Checker / 2026年4月1日時点", 16, RGBColor(213, 219, 225), False)
    add_textbox(slide, Inches(1.15), Inches(5.62), Inches(4.4), Inches(0.35), "2025年7月期　中村隆之", 16, RGBColor(213, 219, 225), False)
    side = add_panel(slide, Inches(9.48), Inches(1.2), Inches(2.95), Inches(4.7), fill_rgb=RGBColor(242, 238, 229), line_rgb=RGBColor(242, 238, 229))
    side.line.fill.background()
    add_textbox(slide, Inches(9.82), Inches(1.62), Inches(2.2), Inches(0.28), "発表で扱う内容", 15, SUBTLE, True)
    add_bullets(
        slide,
        Inches(9.72),
        Inches(2.02),
        Inches(2.35),
        Inches(3.15),
        [
            "課題と選んだ理由",
            "使用したデータ",
            "仮説とアプローチ",
            "分析結果・考察",
            "ビジネスインパクト",
            "工夫点と今後",
        ],
        font_size=15,
        color=INK,
    )

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "1. 課題と選んだ理由", "誤った情報に踊らされず、正確な情報へ近づく支援を作る")
    add_panel(slide, Inches(0.72), Inches(1.52), Inches(5.85), Inches(4.95), fill_rgb=WHITE)
    add_panel(slide, Inches(6.76), Inches(1.52), Inches(5.85), Inches(4.95), fill_rgb=MINT)
    add_textbox(slide, Inches(0.98), Inches(1.82), Inches(2.3), Inches(0.3), "社会的な背景", 20, TEAL, True)
    add_bullets(
        slide,
        Inches(0.95),
        Inches(2.18),
        Inches(5.1),
        Inches(3.6),
        [
            "オールドメディアと呼ばれる既存メディアへの不信感が強まる一方で、X や YouTube など SNS の台頭により、マスメディアとは異なる情報が大量に流通している",
            "2024年の兵庫県知事の辞職と再選という一連の騒動では、一方的な報道しかしないマスメディアへの不信感を強く意識した",
            "世界情勢が不安定な時代だからこそ、戦時中にマスメディアが戦争を煽った歴史も踏まえ、誤った情報に流されない仕組みが必要だと考えた",
        ],
        font_size=17,
    )
    add_textbox(slide, Inches(7.03), Inches(1.82), Inches(2.6), Inches(0.3), "個人的に選んだ理由", 20, TEAL, True)
    add_bullets(
        slide,
        Inches(7.0),
        Inches(2.18),
        Inches(5.1),
        Inches(3.55),
        [
            "昔、誤った情報に騙されて痛い出費をした苦い過去があり、マスメディアに騙されたようで少し怒りを感じた",
            "真偽を自分で確かめる力がなければ、強い言葉や一方的な報道に振り回されてしまうと実感した",
            "だからこそ、断定の代行ではなく、正確な情報へ近づくための確認支援ツールを作りたいと考えた",
        ],
        font_size=17,
    )
    add_metric_card(slide, Inches(6.98), Inches(5.45), Inches(1.55), Inches(0.82), "入力", "URL / 本文", TEAL)
    add_metric_card(slide, Inches(8.76), Inches(5.45), Inches(1.45), Inches(0.82), "出力", "5区分", ORANGE)
    add_metric_card(slide, Inches(10.42), Inches(5.45), Inches(1.88), Inches(0.82), "目標", "正確な情報へ近づく", GREEN)
    add_footer(slide)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "2. 使用したデータについて", "balanced 100件の real_article_dataset_v2 を使って Ver2 / Ver3 を比較評価")
    add_panel(slide, Inches(0.72), Inches(1.62), Inches(3.5), Inches(3.0), fill_rgb=WHITE)
    add_panel(slide, Inches(4.45), Inches(1.62), Inches(3.15), Inches(3.0), fill_rgb=MINT)
    add_panel(slide, Inches(7.83), Inches(1.62), Inches(4.76), Inches(3.0), fill_rgb=WHITE)
    add_textbox(slide, Inches(0.98), Inches(1.92), Inches(2.3), Inches(0.3), "dataset の概要", 18, TEAL, True)
    add_bullets(
        slide,
        Inches(0.95),
        Inches(2.24),
        Inches(2.95),
        Inches(1.95),
        [
            f"ファイル: real_article_dataset_v2.json",
            f"件数: {dataset.case_count}",
            f"curated_on: {dataset.curated_on}",
            "5区分は各20件で balanced",
        ],
        font_size=16,
    )
    add_textbox(slide, Inches(4.72), Inches(1.92), Inches(2.0), Inches(0.3), "スキーマ", 18, TEAL, True)
    add_bullets(
        slide,
        Inches(4.68),
        Inches(2.24),
        Inches(2.4),
        Inches(1.95),
        [f"`{field}`" for field in dataset.fields],
        font_size=16,
    )
    add_textbox(slide, Inches(8.1), Inches(1.92), Inches(2.0), Inches(0.3), "ドメイン分布", 18, TEAL, True)
    add_bullets(
        slide,
        Inches(8.06),
        Inches(2.24),
        Inches(4.0),
        Inches(1.95),
        [
            f"一般: {dataset.domain_counts['一般']}",
            f"医療: {dataset.domain_counts['医療']}",
            f"金融: {dataset.domain_counts['金融']}",
            f"災害: {dataset.domain_counts['災害']}",
            f"政治: {dataset.domain_counts['政治']}",
        ],
        font_size=16,
    )
    add_panel(slide, Inches(0.76), Inches(4.95), Inches(11.82), Inches(1.15), fill_rgb=SAND)
    add_bullets(
        slide,
        Inches(0.98),
        Inches(5.18),
        Inches(11.3),
        Inches(0.6),
        [
            "評価対象は記事全文ではなく analysis_text に要約した中心命題で、expected_verdict と expected_domain を直接比較できるようにしている",
        ],
        font_size=17,
    )
    add_footer(slide)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "3. 初期仮説・分析アプローチ", "ローカル一次判定と Gemini の外部根拠比較を組み合わせれば、境界ケースを改善できると考えた")
    add_panel(slide, Inches(0.74), Inches(1.58), Inches(11.84), Inches(1.82), fill_rgb=SAND)
    add_bullets(
        slide,
        Inches(0.98),
        Inches(1.9),
        Inches(11.25),
        Inches(1.1),
        [
            "仮説1: 出典の有無、著者、日付、リンク数、表現の強さなどのヒューリスティックで危険度の目安は作れる",
            "仮説2: 境界ケースは Gemini の google_search と url_context を使った外部根拠比較を加えることで改善できる",
            "仮説3: 書き振り評価は真偽判定と分離した方が、煽り表現だけで誤判定するリスクを下げられる",
        ],
        font_size=17,
    )
    flow_y = Inches(4.15)
    flow_w = Inches(2.0)
    add_process_box(slide, Inches(0.78), flow_y, flow_w, Inches(1.38), "入力", "URL または本文", TEAL)
    add_connector(slide, Inches(2.9), flow_y + Inches(0.46), Inches(0.35), Inches(0.34), RGBColor(175, 197, 198))
    add_process_box(slide, Inches(3.18), flow_y, flow_w, Inches(1.38), "ローカル一次判定", "risk / confidence / domain を推定", ORANGE)
    add_connector(slide, Inches(5.3), flow_y + Inches(0.46), Inches(0.35), Inches(0.34), RGBColor(220, 192, 173))
    add_process_box(slide, Inches(5.58), flow_y, flow_w, Inches(1.38), "Gemini 比較", "外部根拠と反証を照合", GREEN)
    add_connector(slide, Inches(7.7), flow_y + Inches(0.46), Inches(0.35), Inches(0.34), RGBColor(183, 211, 193))
    add_process_box(slide, Inches(7.98), flow_y, flow_w, Inches(1.38), "書き振り評価", "煽りや断定を別軸で評価", GOLD)
    add_connector(slide, Inches(10.1), flow_y + Inches(0.46), Inches(0.35), Inches(0.34), RGBColor(223, 208, 171))
    add_process_box(slide, Inches(10.38), flow_y, Inches(2.0), Inches(1.38), "出力", "5区分・理由・確認リンク", RED)
    add_footer(slide)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "4. 分析結果", "Ver2 は Ver3 より少し改善したが、誤り recall 0.000 は共通課題として残っている")
    slide.shapes.add_picture(str(plot_path(ver2.plots_dir, "evaluation_dashboard.png")), Inches(0.82), Inches(1.72), width=Inches(6.95))
    add_panel(slide, Inches(8.0), Inches(1.72), Inches(4.28), Inches(4.82), fill_rgb=SAND)
    add_metric_card(slide, Inches(8.28), Inches(2.0), Inches(1.58), Inches(0.82), "Ver2 Accuracy", score(ver2.accuracy), TEAL)
    add_metric_card(slide, Inches(10.08), Inches(2.0), Inches(1.58), Inches(0.82), "Ver2 Macro F1", score(ver2.macro_f1), ORANGE)
    add_metric_card(slide, Inches(8.28), Inches(3.02), Inches(1.58), Inches(0.82), "Ver3 Accuracy", score(ver3.accuracy), TEAL)
    add_metric_card(slide, Inches(10.08), Inches(3.02), Inches(1.58), Inches(0.82), "Ver3 Macro F1", score(ver3.macro_f1), ORANGE)
    add_metric_card(slide, Inches(8.28), Inches(4.04), Inches(1.58), Inches(0.82), "誤り recall", score(ver2.binary_recall), RED)
    add_metric_card(slide, Inches(10.08), Inches(4.04), Inches(1.58), Inches(0.82), "比較", "Ver2優位", GREEN)
    add_bullets(
        slide,
        Inches(8.22),
        Inches(5.12),
        Inches(3.55),
        Inches(1.05),
        [
            "Ver2 は accuracy 0.370、macro F1 0.266",
            "Ver3 は accuracy 0.330、macro F1 0.206",
            "どちらも『誤り』を 1件も拾えていない",
        ],
        font_size=15,
    )
    add_footer(slide)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "4. 考察", "現在の Ver2 は『中間帯の検出』に寄りやすく、ラベルの両端を出し切れていない")
    slide.shapes.add_picture(str(plot_path(ver2.plots_dir, "confusion_matrix.png")), Inches(0.82), Inches(1.78), width=Inches(5.5))
    add_panel(slide, Inches(6.6), Inches(1.76), Inches(5.72), Inches(4.9), fill_rgb=WHITE)
    add_textbox(slide, Inches(6.92), Inches(2.06), Inches(2.6), Inches(0.3), "主な誤判定", 20, ORANGE, True)
    add_bullets(
        slide,
        Inches(6.9),
        Inches(2.38),
        Inches(4.9),
        Inches(1.35),
        [f"{truth} -> {pred}: {count}件" for (truth, pred), count in ver2.mismatch_counts[:3]],
        font_size=18,
    )
    add_textbox(slide, Inches(6.92), Inches(4.05), Inches(2.6), Inches(0.3), "考察", 20, TEAL, True)
    add_bullets(
        slide,
        Inches(6.9),
        Inches(4.38),
        Inches(4.9),
        Inches(1.65),
        [
            "正確を『正確』へ戻せず、すべて『ほぼ正確』へ寄せている",
            "誤りを『判断保留』や『不正確』へ逃がし、危険側へ押し切れていない",
            "直近 tune4 の境界調整は改善せず、evidence の重み付け見直しが必要",
        ],
        font_size=18,
    )
    add_footer(slide)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "5. ビジネスインパクト", "現時点では『自動判定ツール』より『説明付き確認支援ツール』として価値が高い")
    add_metric_card(slide, Inches(0.88), Inches(1.82), Inches(3.55), Inches(1.02), "個人向け", "投稿前セルフチェック", TEAL)
    add_metric_card(slide, Inches(4.88), Inches(1.82), Inches(3.55), Inches(1.02), "組織向け", "一次仕分けと優先順位付け", ORANGE)
    add_metric_card(slide, Inches(8.88), Inches(1.82), Inches(3.48), Inches(1.02), "教育向け", "メディアリテラシー支援", GREEN)
    add_panel(slide, Inches(0.88), Inches(3.2), Inches(11.45), Inches(3.0), fill_rgb=SAND)
    add_bullets(
        slide,
        Inches(1.15),
        Inches(3.5),
        Inches(10.9),
        Inches(2.2),
        [
            "5区分 verdict に加え、理由・確認リンク・attention score を返すので、利用者が次の確認行動を取りやすい",
            "自治体、学校、メディア運用では『まず人が見るべき候補』を絞る補助として使える",
            "ただし 2026-04-01 時点では誤り recall が 0.000 のため、単独の自動ジャッジ用途にはまだ不十分",
            "価値の中心は、断定の代行ではなく、確認コストの削減と注意喚起にある",
        ],
        font_size=18,
    )
    add_footer(slide)

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "6. 工夫点", "評価可能な形に落とし込み、レビューしやすい運用まで整えた")
    add_panel(slide, Inches(0.74), Inches(1.62), Inches(5.9), Inches(4.95), fill_rgb=WHITE)
    add_panel(slide, Inches(6.86), Inches(1.62), Inches(5.7), Inches(4.95), fill_rgb=MINT)
    add_textbox(slide, Inches(1.0), Inches(1.92), Inches(2.2), Inches(0.3), "設計面の工夫", 20, TEAL, True)
    add_bullets(
        slide,
        Inches(0.98),
        Inches(2.26),
        Inches(5.1),
        Inches(2.15),
        [
            "真偽判定と書き振り評価を分離し、煽り表現だけで真偽を決めない設計にした",
            "dataset_runner / evaluation / plot をつなぎ、同じ dataset で改善前後を比較できるようにした",
            "Gemini が失敗してもローカル判定へフォールバックして止まらないようにした",
        ],
        font_size=17,
    )
    add_textbox(slide, Inches(1.0), Inches(4.78), Inches(2.8), Inches(0.3), "運用面の工夫", 20, ORANGE, True)
    add_bullets(
        slide,
        Inches(0.98),
        Inches(5.1),
        Inches(5.05),
        Inches(1.0),
        [
            "evaluation_outputs では 5 plots を必須で出す運用にした",
            "predictions JSON / eval JSON / plots / CSV / XLSX を同じ timestamp フォルダへまとめた",
            "export_evaluation_bundle.ps1 で評価一式をまとめて再生成できるようにした",
        ],
        font_size=17,
    )
    positions = [
        (Inches(7.08), Inches(1.95), Inches(2.2), Inches(1.55), "evaluation_dashboard.png"),
        (Inches(9.46), Inches(1.95), Inches(2.2), Inches(1.55), "confusion_matrix.png"),
        (Inches(7.08), Inches(3.72), Inches(2.2), Inches(1.55), "summary_metrics.png"),
        (Inches(9.46), Inches(3.72), Inches(2.2), Inches(1.55), "per_class_metrics.png"),
    ]
    for left, top, width, height, filename in positions:
        add_labeled_image(slide, plot_path(ver2.plots_dir, filename), left, top, width, height, filename)
    add_textbox(slide, Inches(7.1), Inches(5.5), Inches(4.2), Inches(0.26), "必須 5 枚目の evaluation_overview.png も同じ bundle に保存", 12, SUBTLE, True)
    add_footer(slide)

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "6. 今後さらに実施してみたいこと", "Ver2 を中心に、誤り経路の強化と正確ラベル復帰を次の改善テーマにする")
    add_panel(slide, Inches(0.76), Inches(1.72), Inches(5.82), Inches(4.92), fill_rgb=WHITE)
    add_panel(slide, Inches(6.82), Inches(1.72), Inches(5.82), Inches(4.92), fill_rgb=MINT)
    add_textbox(slide, Inches(1.02), Inches(2.02), Inches(2.8), Inches(0.3), "モデル改善", 20, TEAL, True)
    add_bullets(
        slide,
        Inches(1.0),
        Inches(2.34),
        Inches(5.05),
        Inches(3.4),
        [
            "反証ありケースから『誤り』へ落とす evidence 設計を見直す",
            "正確なケースを『ほぼ正確』へ寄せすぎないよう、正確ラベル復帰の較正を進める",
            "tune4 を基準比較し、必要なら一段前の閾値へ戻して差分を測る",
        ],
        font_size=17,
    )
    add_textbox(slide, Inches(7.08), Inches(2.02), Inches(3.6), Inches(0.3), "データと運用の拡張", 20, ORANGE, True)
    add_bullets(
        slide,
        Inches(7.05),
        Inches(2.34),
        Inches(5.05),
        Inches(3.4),
        [
            "CSV / XLSX を使って誤判定をケース単位で見直し、ルール修正の根拠を増やす",
            "時事性の高いケースや SNS 投稿に近いデータを追加し、実運用に近い評価へ広げる",
            "発表では『今できていること』と『まだできていないこと』を分けて伝える",
        ],
        font_size=17,
    )
    add_footer(slide)

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "図の説明: summary_metrics.png", "全体指標を見ると、Ver2 はまだ『誤り』検出で止まっている")
    slide.shapes.add_picture(str(plot_path(ver2.plots_dir, "summary_metrics.png")), Inches(0.82), Inches(1.72), width=Inches(6.9))
    add_panel(slide, Inches(8.0), Inches(1.72), Inches(4.3), Inches(4.85), fill_rgb=SAND)
    add_textbox(slide, Inches(8.26), Inches(2.0), Inches(2.8), Inches(0.3), "読み取りポイント", 19, TEAL, True)
    add_bullets(
        slide,
        Inches(8.22),
        Inches(2.34),
        Inches(3.6),
        Inches(2.8),
        [
            "Accuracy は 0.370、Macro F1 は 0.266 にとどまっている",
            "Weighted F1 も同じ 0.266 で、全体として強いクラス優位は作れていない",
            "Binary Precision / Recall / F1 がすべて 0.000 で、誤りラベルを 1 件も当てられていない",
            "つまり全体の課題は『誤情報を危険側へ押し切れないこと』に集約される",
        ],
        font_size=17,
    )
    add_footer(slide)

    # Slide 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "図の説明: per_class_metrics.png", "クラスごとに見ると、両端ラベルが崩れて中間へ寄っている")
    slide.shapes.add_picture(str(plot_path(ver2.plots_dir, "per_class_metrics.png")), Inches(0.72), Inches(1.72), width=Inches(7.45))
    add_panel(slide, Inches(8.35), Inches(1.72), Inches(3.95), Inches(4.85), fill_rgb=WHITE)
    add_textbox(slide, Inches(8.62), Inches(2.0), Inches(2.8), Inches(0.3), "読み取りポイント", 19, TEAL, True)
    add_bullets(
        slide,
        Inches(8.58),
        Inches(2.34),
        Inches(3.2),
        Inches(2.85),
        [
            "『正確』と『誤り』は Precision / Recall / F1 がすべて 0.000",
            "『ほぼ正確』は Recall 0.900 と高く、正確側の多くをここへ吸収している",
            "『不正確』は Recall 0.650 で、中間より危険側はある程度拾えている",
            "結果として、モデルは白黒をつけず『ほぼ正確』『判断保留』『不正確』へ寄せやすい",
        ],
        font_size=17,
    )
    add_footer(slide)

    # Slide 12
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "図の説明: confusion_matrix.png", "どのラベルへ流れているかを見ると、誤りを『判断保留』や『不正確』へ逃がしている")
    slide.shapes.add_picture(str(plot_path(ver2.plots_dir, "confusion_matrix.png")), Inches(0.82), Inches(1.7), width=Inches(6.55))
    add_panel(slide, Inches(7.75), Inches(1.72), Inches(4.55), Inches(4.85), fill_rgb=MINT)
    add_textbox(slide, Inches(8.02), Inches(2.0), Inches(3.1), Inches(0.3), "読み取りポイント", 19, TEAL, True)
    add_bullets(
        slide,
        Inches(7.98),
        Inches(2.34),
        Inches(3.7),
        Inches(2.85),
        [
            "正確 20 件はすべて『ほぼ正確』へ流れている",
            "誤り 20 件は『判断保留』9 件、『不正確』11 件へ分散している",
            "判断保留 20 件のうち 12 件が『不正確』へ寄っており、灰色領域の扱いも不安定",
            "この図から、閾値だけでなく evidence の使い方自体を見直す必要があると分かる",
        ],
        font_size=17,
    )
    add_footer(slide)

    # Slide 13
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "図の説明: evaluation_overview.png", "集計サマリーからも、mismatch の多さと誤り検出の弱さが確認できる")
    slide.shapes.add_picture(str(plot_path(ver2.plots_dir, "evaluation_overview.png")), Inches(1.2), Inches(1.68), width=Inches(4.6))
    add_panel(slide, Inches(6.35), Inches(1.72), Inches(5.95), Inches(4.85), fill_rgb=SAND)
    add_textbox(slide, Inches(6.65), Inches(2.0), Inches(3.0), Inches(0.3), "読み取りポイント", 19, TEAL, True)
    add_bullets(
        slide,
        Inches(6.6),
        Inches(2.34),
        Inches(5.1),
        Inches(2.75),
        [
            "sample_count は 100、skipped_count は 0 なので、全件評価はできている",
            "mismatches は 63 件あり、正解より不一致の方がかなり多い",
            "Macro Avg precision 0.2131 / recall 0.3700 / f1 0.2659 で、総合性能はまだ改善余地が大きい",
            "positive_label を『誤り』に置いた評価でも support 20 に対して成果が出ておらず、ここが次の最優先課題",
        ],
        font_size=17,
    )
    add_footer(slide)

    # Slide 14
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "参考資料: confusion_matrix.png", "Ver2 / evaluation_outputs/20260401-1213/plots")
    slide.shapes.add_picture(str(plot_path(ver2.plots_dir, "confusion_matrix.png")), Inches(1.05), Inches(1.55), width=Inches(11.1))
    add_footer(slide)

    # Slide 15
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "参考資料: evaluation_overview.png", "Ver2 / evaluation_outputs/20260401-1213/plots")
    slide.shapes.add_picture(str(plot_path(ver2.plots_dir, "evaluation_overview.png")), Inches(3.35), Inches(1.45), width=Inches(6.5))
    add_footer(slide)

    # Slide 16
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "参考資料: per_class_metrics.png", "Ver2 / evaluation_outputs/20260401-1213/plots")
    slide.shapes.add_picture(str(plot_path(ver2.plots_dir, "per_class_metrics.png")), Inches(0.85), Inches(1.6), width=Inches(11.6))
    add_footer(slide)

    # Slide 17
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header(slide, "参考資料: summary_metrics.png", "Ver2 / evaluation_outputs/20260401-1213/plots")
    slide.shapes.add_picture(str(plot_path(ver2.plots_dir, "summary_metrics.png")), Inches(0.9), Inches(1.65), width=Inches(11.5))
    add_footer(slide)

    return prs


def main() -> None:
    PRESENTATION_DIR.mkdir(parents=True, exist_ok=True)
    dataset = build_dataset_summary(DATASET_PATH)
    ver2 = build_eval_snapshot("Ver2 tune4", VER2_EVAL_PATH, VER2_PREDICTIONS_PATH, VER2_PLOTS_DIR)
    ver3 = build_eval_snapshot("Ver3 tune3", VER3_EVAL_PATH, VER3_PREDICTIONS_PATH, VER3_PLOTS_DIR)

    prs = build_presentation(dataset, ver2, ver3)
    output_pptx = OUTPUT_PPTX
    try:
        prs.save(output_pptx)
    except PermissionError:
        output_pptx = FALLBACK_OUTPUT_PPTX
        prs.save(output_pptx)
    OUTPUT_NOTES.write_text(build_notes(dataset, ver2, ver3), encoding="utf-8")

    print(output_pptx)
    print(OUTPUT_NOTES)


if __name__ == "__main__":
    main()
