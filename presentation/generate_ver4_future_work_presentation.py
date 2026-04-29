from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PRESENTATION_DIR = Path(__file__).resolve().parent
PROJECT_ROOT = PRESENTATION_DIR.parent
EVAL_PATH = PROJECT_ROOT / "Ver4" / "evaluation_outputs" / "20260427-0304" / "eval_real_article_dataset_v2_use_gpt_gemini.json"

OUTPUT_PPTX = PRESENTATION_DIR / "ver4_future_work_only_20260429.pptx"
OUTPUT_NOTES = PRESENTATION_DIR / "ver4_future_work_only_20260429_notes.md"

FONT = "Yu Gothic"
BG = RGBColor(247, 248, 246)
INK = RGBColor(27, 35, 43)
SUBTLE = RGBColor(88, 99, 110)
WHITE = RGBColor(255, 255, 255)
NAVY = RGBColor(24, 48, 72)
TEAL = RGBColor(15, 124, 128)
GREEN = RGBColor(61, 143, 99)
GOLD = RGBColor(191, 151, 48)
ORANGE = RGBColor(207, 115, 47)
RED = RGBColor(182, 63, 63)
LINE = RGBColor(214, 220, 224)
MINT = RGBColor(230, 243, 240)
BLUE_SOFT = RGBColor(232, 239, 248)
GREEN_SOFT = RGBColor(232, 244, 235)
GOLD_SOFT = RGBColor(251, 246, 226)
ORANGE_SOFT = RGBColor(250, 238, 227)
RED_SOFT = RGBColor(249, 232, 232)


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def pct(value: float) -> str:
    return f"{value * 100:.1f}%"


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_text(slide, left, top, width, height, text="", size=20, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(3)
    tf.margin_right = Pt(3)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if text:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        p.alignment = align
    return box


def add_header(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, Inches(0.62), Inches(0.34), Inches(11.0), Inches(0.42), title, 25, NAVY, True)
    if subtitle:
        add_text(slide, Inches(0.64), Inches(0.81), Inches(11.2), Inches(0.26), subtitle, 10, SUBTLE)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(1.1), Inches(12.05), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def add_footer(slide, page: int) -> None:
    add_text(slide, Inches(0.65), Inches(7.05), Inches(7.4), Inches(0.18), "Fake News Checker Ver4 | 6. 今後さらに実施してみたいこと", 8, SUBTLE)
    add_text(slide, Inches(12.25), Inches(7.05), Inches(0.55), Inches(0.18), str(page), 8, SUBTLE, False, PP_ALIGN.RIGHT)


def add_panel(slide, left, top, width, height, fill=WHITE, line=LINE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_bullets(slide, left, top, width, height, items: list[str], size=16, color=INK):
    box = add_text(slide, left, top, width, height)
    tf = box.text_frame
    tf.clear()
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(5)
        p.line_spacing = 1.1


def add_metric(slide, left, top, width, height, label: str, value: str, accent=TEAL, note: str | None = None):
    add_panel(slide, left, top, width, height)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.08))
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()
    add_text(slide, left + Inches(0.12), top + Inches(0.2), width - Inches(0.24), Inches(0.2), label, 11, SUBTLE, True, PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.12), top + Inches(0.47), width - Inches(0.24), Inches(0.42), value, 24, accent, True, PP_ALIGN.CENTER)
    if note:
        add_text(slide, left + Inches(0.12), top + Inches(0.93), width - Inches(0.24), Inches(0.23), note, 9, SUBTLE, False, PP_ALIGN.CENTER)


def add_card(slide, left, top, width, height, title: str, body: str, accent, fill):
    add_panel(slide, left, top, width, height, fill)
    add_text(slide, left + Inches(0.22), top + Inches(0.22), width - Inches(0.44), Inches(0.25), title, 17, accent, True)
    add_text(slide, left + Inches(0.22), top + Inches(0.64), width - Inches(0.44), height - Inches(0.78), body, 12, INK)


def add_step(slide, left, top, width, height, number: str, title: str, body: str, accent):
    add_panel(slide, left, top, width, height, WHITE)
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.2), top + Inches(0.18), Inches(0.4), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = accent
    badge.line.fill.background()
    add_text(slide, left + Inches(0.2), top + Inches(0.25), Inches(0.4), Inches(0.15), number, 9, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.82), top + Inches(0.2), Inches(3.5), Inches(0.22), title, 15, accent, True)
    add_text(slide, left + Inches(4.48), top + Inches(0.22), width - Inches(4.75), Inches(0.2), body, 12, INK)


def build_notes(eval_report: dict) -> str:
    return "\n".join(
        [
            "# Ver4 Future Work Notes",
            "",
            "- generated_on: 2026-04-29",
            f"- current_accuracy: {pct(eval_report['multiclass']['accuracy'])}",
            f"- current_macro_f1: {pct(eval_report['multiclass']['macro_avg']['f1'])}",
            "",
            "## Slide 1",
            "- 今後さらに実施してみたいことのみの抜粋資料。",
            "",
            "## Slide 2",
            "- 既存100件では高性能だが、改善に使ったデータでもあるため、未知データでの確認が必要。",
            "",
            "## Slide 3",
            "- blind test 100件を作り、Ver2/3/4/5を同じ条件で比較する。",
            "",
            "## Slide 4",
            "- 境界判定ルールを固有名詞パッチではなく、一般化できる構造ルールにする。",
            "",
            "## Slide 5",
            "- 実運用を想定し、処理時間、根拠リンク、UI、一括処理、修正履歴の改善を行う。",
            "",
            "## Slide 6",
            "- devとblindの両方で確認し、既存性能と汎化性能のバランスを見る。",
            "",
            "## Slide 7",
            "- まとめとして、次の焦点は未知データで本当に使えるかの検証だと伝える。",
        ]
    ) + "\n"


def build_presentation() -> None:
    eval_report = load_json(EVAL_PATH)
    accuracy = eval_report["multiclass"]["accuracy"]
    macro_f1 = eval_report["multiclass"]["macro_avg"]["f1"]
    mismatch_count = len(eval_report["mismatches"])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    add_text(slide, Inches(0.82), Inches(1.22), Inches(10), Inches(0.42), "Fake News Checker Ver4", 25, RGBColor(205, 229, 227), True)
    add_text(slide, Inches(0.82), Inches(1.78), Inches(11.6), Inches(1.05), "6. 今後さらに実施してみたいこと", 40, WHITE, True)
    add_text(slide, Inches(0.86), Inches(3.02), Inches(10.0), Inches(0.55), "既存100件での高性能を、未知データで確認する", 20, RGBColor(232, 238, 244))
    add_metric(slide, Inches(0.9), Inches(4.55), Inches(2.45), Inches(1.15), "Current Accuracy", pct(accuracy), TEAL)
    add_metric(slide, Inches(3.55), Inches(4.55), Inches(2.45), Inches(1.15), "Macro F1", pct(macro_f1), GREEN)
    add_metric(slide, Inches(6.2), Inches(4.55), Inches(2.45), Inches(1.15), "Mismatches", str(mismatch_count), ORANGE)
    add_text(slide, Inches(0.9), Inches(6.38), Inches(4.3), Inches(0.22), "2026-04-29", 12, RGBColor(214, 224, 232))

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "なぜ次の検証が必要か", "既存100件では高性能。ただし、未知データでの確認がまだ必要")
    add_panel(slide, Inches(0.82), Inches(1.48), Inches(5.65), Inches(3.15), GREEN_SOFT)
    add_text(slide, Inches(1.12), Inches(1.83), Inches(5.0), Inches(0.3), "現在できていること", 21, GREEN, True)
    add_bullets(
        slide,
        Inches(1.1),
        Inches(2.42),
        Inches(5.0),
        Inches(1.45),
        [
            f"既存100件でAccuracy {pct(accuracy)}",
            f"Macro F1 {pct(macro_f1)}まで改善",
            "5段階分類でも一定の精度を確認",
        ],
        15,
    )
    add_panel(slide, Inches(6.82), Inches(1.48), Inches(5.65), Inches(3.15), ORANGE_SOFT)
    add_text(slide, Inches(7.12), Inches(1.83), Inches(5.0), Inches(0.3), "まだ確認したいこと", 21, ORANGE, True)
    add_bullets(
        slide,
        Inches(7.1),
        Inches(2.42),
        Inches(5.0),
        Inches(1.45),
        [
            "既存100件に合わせすぎていないか",
            "新しい記事や主張でも性能が出るか",
            "中間ラベルの境界が安定するか",
        ],
        15,
    )
    add_text(slide, Inches(0.92), Inches(5.82), Inches(11.45), Inches(0.38), "次の焦点: 既存データで良いだけでなく、未知データでも通用するかを確認する。", 18, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 2)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "検証計画", "blind testを作り、各バージョンを同じ条件で比較する")
    steps = [
        ("1", "blind test 100件を作る", "改善に使っていない新しい評価データを用意する", TEAL),
        ("2", "Ver2/3/4/5を同じblind setで評価", "方式の違いを公平に比較する", GREEN),
        ("3", "devとblindの両方で確認", "既存性能を落とさず、汎化性能を上げる", GOLD),
        ("4", "ミスマッチを再分析", "どの境界で失敗するかを構造的に見る", ORANGE),
    ]
    for i, (num, title, body, color) in enumerate(steps):
        add_step(slide, Inches(0.95), Inches(1.48 + i * 1.05), Inches(11.45), Inches(0.72), num, title, body, color)
    add_panel(slide, Inches(1.08), Inches(5.88), Inches(11.15), Inches(0.62), BLUE_SOFT)
    add_text(slide, Inches(1.35), Inches(6.08), Inches(10.6), Inches(0.18), "dev = 改善に使う既存データ、blind = 改善には使わない新規検証データ", 14, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 3)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "境界判定ルールの改善", "固有名詞ごとの個別対応ではなく、一般化できるルールにする")
    rules = [
        ("正確 vs ほぼ正確", "数字・時期・表現のズレが、結論に影響するかを見る。", TEAL, MINT),
        ("ほぼ正確 vs 不正確", "一部の誤りが軽微か、主張の印象を変える重要な誤りかを分ける。", ORANGE, ORANGE_SOFT),
        ("不正確 vs 誤り", "一部真実があるのか、核心が明確に否定されるのかを分ける。", RED, RED_SOFT),
        ("判断保留 vs 不正確", "根拠不足なのか、否定根拠が十分あるのかを分ける。", GOLD, GOLD_SOFT),
    ]
    for i, (title, body, color, fill) in enumerate(rules):
        left = Inches(0.78 + (i % 2) * 6.05)
        top = Inches(1.55 + (i // 2) * 1.78)
        add_card(slide, left, top, Inches(5.55), Inches(1.32), title, body, color, fill)
    add_panel(slide, Inches(1.08), Inches(5.58), Inches(11.15), Inches(0.82), BLUE_SOFT)
    add_text(slide, Inches(1.35), Inches(5.86), Inches(10.6), Inches(0.25), "目標: 個別ケースへのパッチではなく、未知データにも効く判定基準へ整理する。", 17, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 4)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "実運用を想定した改善", "精度だけでなく、使いやすさと運用コストも改善する")
    improvements = [
        ("処理時間の短縮", "大量記事を扱うため、判定速度とAPI利用コストを改善する。", TEAL, MINT),
        ("根拠リンクの信頼性評価", "公式情報・報道・SNSなど、根拠の信頼度を区別する。", GREEN, GREEN_SOFT),
        ("UI表示の改善", "判定、理由、注意点、確認リンクを見やすく整理する。", GOLD, GOLD_SOFT),
        ("一括処理と記録化", "複数記事のCSV出力や、人間の修正履歴を残せるようにする。", ORANGE, ORANGE_SOFT),
    ]
    for i, (title, body, color, fill) in enumerate(improvements):
        left = Inches(0.78 + (i % 2) * 6.05)
        top = Inches(1.55 + (i // 2) * 1.78)
        add_card(slide, left, top, Inches(5.55), Inches(1.32), title, body, color, fill)
    add_text(slide, Inches(0.92), Inches(6.28), Inches(11.45), Inches(0.28), "実運用では、精度に加えて、速さ・説明の見やすさ・再確認のしやすさが重要。", 15, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 5)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "評価の見方", "devとblindの両方で確認する")
    add_panel(slide, Inches(0.95), Inches(1.55), Inches(5.4), Inches(3.25), MINT)
    add_text(slide, Inches(1.24), Inches(1.92), Inches(4.85), Inches(0.34), "devで確認すること", 22, TEAL, True)
    add_bullets(
        slide,
        Inches(1.22),
        Inches(2.56),
        Inches(4.8),
        Inches(1.25),
        [
            "既存100件で性能が落ちていないか",
            "改善前に直した誤分類が再発していないか",
            "ルール変更の副作用がないか",
        ],
        14,
    )
    add_panel(slide, Inches(6.95), Inches(1.55), Inches(5.4), Inches(3.25), ORANGE_SOFT)
    add_text(slide, Inches(7.24), Inches(1.92), Inches(4.85), Inches(0.34), "blindで確認すること", 22, ORANGE, True)
    add_bullets(
        slide,
        Inches(7.22),
        Inches(2.56),
        Inches(4.8),
        Inches(1.25),
        [
            "未知データにも対応できるか",
            "既存100件に合わせすぎていないか",
            "汎化性能があるか",
        ],
        14,
    )
    add_panel(slide, Inches(1.08), Inches(5.58), Inches(11.15), Inches(0.82), BLUE_SOFT)
    add_text(slide, Inches(1.35), Inches(5.86), Inches(10.6), Inches(0.25), "両方で見ることで、既存性能と汎化性能のバランスを確認する。", 18, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 6)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "まとめ", "次は、未知データで本当に使えるかを確認する")
    add_panel(slide, Inches(1.05), Inches(1.58), Inches(11.1), Inches(1.2), MINT)
    add_text(slide, Inches(1.32), Inches(1.92), Inches(10.58), Inches(0.34), "既存100件で高性能 → blind testで汎化性能を検証", 24, NAVY, True, PP_ALIGN.CENTER)
    add_bullets(
        slide,
        Inches(1.45),
        Inches(3.35),
        Inches(10.35),
        Inches(1.85),
        [
            "新規blind 100件を作り、既存データへの合わせすぎを確認する",
            "Ver2/3/4/5を同じ条件で比較し、方式差を公平に見る",
            "中間ラベルの境界判定ルールを、一般化できる形で改善する",
            "実運用を見据えて、速度・UI・根拠リンク・記録化も改善する",
        ],
        17,
    )
    add_text(slide, Inches(1.15), Inches(6.14), Inches(11.0), Inches(0.38), "プレゼンでの一言: 今後は、精度を上げるだけでなく、未知データでも安定して使えるかを検証します。", 16, TEAL, True, PP_ALIGN.CENTER)
    add_footer(slide, 7)

    prs.save(OUTPUT_PPTX)
    OUTPUT_NOTES.write_text(build_notes(eval_report), encoding="utf-8")


if __name__ == "__main__":
    build_presentation()
    print(f"saved: {OUTPUT_PPTX}")
    print(f"saved: {OUTPUT_NOTES}")
