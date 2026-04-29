from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PRESENTATION_DIR = Path(__file__).resolve().parent
OUTPUT_PPTX = PRESENTATION_DIR / "ver4_business_impact_only_20260429.pptx"
OUTPUT_NOTES = PRESENTATION_DIR / "ver4_business_impact_only_20260429_notes.md"

FONT = "Yu Gothic"
BG = RGBColor(247, 248, 246)
INK = RGBColor(27, 35, 43)
SUBTLE = RGBColor(88, 99, 110)
WHITE = RGBColor(255, 255, 255)
NAVY = RGBColor(24, 48, 72)
TEAL = RGBColor(15, 124, 128)
GREEN = RGBColor(61, 143, 99)
GOLD = RGBColor(191, 151, 48)
ORANGE = RGBColor(207, 115, 47)
RED = RGBColor(182, 63, 63)
LINE = RGBColor(214, 220, 224)
MINT = RGBColor(230, 243, 240)
BLUE_SOFT = RGBColor(232, 239, 248)
GREEN_SOFT = RGBColor(232, 244, 235)
GOLD_SOFT = RGBColor(251, 246, 226)
ORANGE_SOFT = RGBColor(250, 238, 227)
RED_SOFT = RGBColor(249, 232, 232)


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_text(slide, left, top, width, height, text="", size=20, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(3)
    tf.margin_right = Pt(3)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if text:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        p.alignment = align
    return box


def add_header(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, Inches(0.62), Inches(0.34), Inches(10.9), Inches(0.42), title, 25, NAVY, True)
    if subtitle:
        add_text(slide, Inches(0.64), Inches(0.81), Inches(11.2), Inches(0.26), subtitle, 10, SUBTLE)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(1.1), Inches(12.05), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def add_footer(slide, page: int) -> None:
    add_text(slide, Inches(0.65), Inches(7.05), Inches(7.2), Inches(0.18), "Fake News Checker Ver4 | 5. ビジネスインパクト", 8, SUBTLE)
    add_text(slide, Inches(12.25), Inches(7.05), Inches(0.55), Inches(0.18), str(page), 8, SUBTLE, False, PP_ALIGN.RIGHT)


def add_panel(slide, left, top, width, height, fill=WHITE, line=LINE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_bullets(slide, left, top, width, height, items: list[str], size=16, color=INK):
    box = add_text(slide, left, top, width, height)
    tf = box.text_frame
    tf.clear()
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(5)
        p.line_spacing = 1.1


def add_impact_card(slide, left, top, width, height, title: str, body: str, accent, fill):
    add_panel(slide, left, top, width, height, fill)
    add_text(slide, left + Inches(0.22), top + Inches(0.22), width - Inches(0.44), Inches(0.25), title, 17, accent, True)
    add_text(slide, left + Inches(0.22), top + Inches(0.64), width - Inches(0.44), height - Inches(0.76), body, 12, INK)


def add_step(slide, left, top, width, height, number: str, title: str, body: str, accent):
    add_panel(slide, left, top, width, height, WHITE)
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.2), top + Inches(0.2), Inches(0.42), Inches(0.42))
    badge.fill.solid()
    badge.fill.fore_color.rgb = accent
    badge.line.fill.background()
    add_text(slide, left + Inches(0.2), top + Inches(0.27), Inches(0.42), Inches(0.16), number, 10, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.85), top + Inches(0.22), Inches(3.4), Inches(0.22), title, 16, accent, True)
    add_text(slide, left + Inches(4.45), top + Inches(0.24), width - Inches(4.75), Inches(0.2), body, 12, INK)


def add_arrow(slide, left, top, width, height, color=SUBTLE):
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()


def build_notes() -> str:
    return "\n".join(
        [
            "# Ver4 Business Impact Notes",
            "",
            "- generated_on: 2026-04-29",
            "",
            "## Slide 1",
            "- ビジネスインパクトのみの抜粋資料。",
            "- Ver4は人間の最終判断を置き換えるのではなく、確認作業を支援するツールとして説明する。",
            "",
            "## Slide 2",
            "- 大量のニュース記事やSNS投稿を、すべて人手で確認するのは難しい。",
            "- Ver4は要注意情報を優先的に見つける一次スクリーニングとして使える。",
            "",
            "## Slide 3",
            "- 主な価値は、一次スクリーニング、調査時間の短縮、説明責任の補助、早期リスク対応。",
            "",
            "## Slide 4",
            "- 業務フローでは、収集、AI一次判定、人間の確認、対応判断、記録化の流れで使う。",
            "",
            "## Slide 5",
            "- 活用場面は、メディア、企業広報、自治体、SNS監視、金融・医療など。",
            "",
            "## Slide 6",
            "- まとめとして、Ver4は根拠付きで人間の確認を速くする支援ツールであると伝える。",
        ]
    ) + "\n"


def build_presentation() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    add_text(slide, Inches(0.82), Inches(1.28), Inches(10), Inches(0.42), "Fake News Checker Ver4", 25, RGBColor(205, 229, 227), True)
    add_text(slide, Inches(0.82), Inches(1.88), Inches(10.9), Inches(1.05), "5. ビジネスインパクト", 46, WHITE, True)
    add_text(slide, Inches(0.86), Inches(3.08), Inches(9.4), Inches(0.55), "大量情報の確認作業を、速く・根拠付きにする", 20, RGBColor(232, 238, 244))
    add_panel(slide, Inches(0.88), Inches(4.55), Inches(10.65), Inches(0.84), RGBColor(37, 67, 96), RGBColor(37, 67, 96))
    add_text(slide, Inches(1.12), Inches(4.84), Inches(10.2), Inches(0.23), "AIが最終判断するのではなく、人間の確認作業を支援する", 18, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(0.9), Inches(6.38), Inches(4.3), Inches(0.22), "2026-04-29", 12, RGBColor(214, 224, 232))

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "業務上の課題", "誤情報は拡散が速く、人手確認だけでは追いつきにくい")
    add_panel(slide, Inches(0.82), Inches(1.55), Inches(5.65), Inches(3.18), BLUE_SOFT)
    add_text(slide, Inches(1.12), Inches(1.9), Inches(5.0), Inches(0.3), "現場で起きること", 21, NAVY, True)
    add_bullets(
        slide,
        Inches(1.1),
        Inches(2.48),
        Inches(5.0),
        Inches(1.5),
        [
            "SNSやWeb上に大量の情報が流れる",
            "誤情報は短時間で拡散しやすい",
            "人手だけでは確認の優先順位をつけにくい",
        ],
        15,
    )
    add_panel(slide, Inches(6.82), Inches(1.55), Inches(5.65), Inches(3.18), MINT)
    add_text(slide, Inches(7.12), Inches(1.9), Inches(5.0), Inches(0.3), "Ver4でできること", 21, TEAL, True)
    add_bullets(
        slide,
        Inches(7.1),
        Inches(2.48),
        Inches(5.0),
        Inches(1.5),
        [
            "要注意の投稿や記事を一次抽出する",
            "主張・理由・確認リンクを整理する",
            "担当者が確認すべき対象を絞り込む",
        ],
        15,
    )
    add_text(slide, Inches(0.92), Inches(5.85), Inches(11.45), Inches(0.36), "価値: すべてを同じ優先度で読む状態から、確認すべきものから見る状態へ変える。", 18, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 2)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "主なビジネスインパクト", "確認作業の効率化と、説明しやすい対応につながる")
    impacts = [
        ("一次スクリーニング", "大量の記事・投稿から、要注意度の高いものを優先的に抽出する。", TEAL, MINT),
        ("調査時間の短縮", "判定理由や確認リンクを出すため、ゼロから調べる時間を減らせる。", GREEN, GREEN_SOFT),
        ("説明責任の補助", "なぜその判定になったかを残せるため、報告や注意喚起に使いやすい。", GOLD, GOLD_SOFT),
        ("早期リスク対応", "拡散前・炎上前の段階で、広報や危機管理の初動を早められる。", ORANGE, ORANGE_SOFT),
    ]
    for i, (title, body, color, fill) in enumerate(impacts):
        left = Inches(0.78 + (i % 2) * 6.05)
        top = Inches(1.55 + (i // 2) * 1.78)
        add_impact_card(slide, left, top, Inches(5.55), Inches(1.32), title, body, color, fill)
    add_panel(slide, Inches(1.08), Inches(5.58), Inches(11.15), Inches(0.82), BLUE_SOFT)
    add_text(slide, Inches(1.35), Inches(5.86), Inches(10.6), Inches(0.25), "最終判断は人間が行い、Ver4は確認対象と根拠を整理する役割を持つ。", 17, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 3)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "業務フローでの使い方", "人間の確認プロセスに組み込む")
    flow = [
        ("1", "情報収集", "SNS投稿・ニュース記事・問い合わせ内容を収集", TEAL),
        ("2", "AI一次判定", "Ver4で要注意度、判定、理由、確認リンクを出力", GREEN),
        ("3", "人間の確認", "担当者が根拠を確認し、最終判断を行う", GOLD),
        ("4", "対応判断", "訂正、注意喚起、社内共有、追加調査を決める", ORANGE),
        ("5", "記録化", "判定理由と根拠をレポートとして残す", RED),
    ]
    for i, (num, title, body, color) in enumerate(flow):
        top = Inches(1.38 + i * 0.92)
        add_step(slide, Inches(0.95), top, Inches(11.45), Inches(0.66), num, title, body, color)
        if i < len(flow) - 1:
            add_arrow(slide, Inches(6.2), top + Inches(0.68), Inches(0.25), Inches(0.18), SUBTLE)
    add_text(slide, Inches(0.92), Inches(6.35), Inches(11.45), Inches(0.28), "導入イメージ: 既存の監視・広報・調査フローの前段に、AIの一次確認を入れる。", 15, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 4)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "活用できる場面", "誤情報が判断や行動に影響する領域で価値が出やすい")
    uses = [
        ("メディア", "記事やSNS上の主張を確認し、ファクトチェックの初動を速くする。", TEAL, MINT),
        ("企業広報・危機管理", "自社や業界に関する誤情報を早期に検知し、対応を判断する。", GREEN, GREEN_SOFT),
        ("自治体・公共機関", "災害時や緊急時の不確かな情報を確認し、住民向け発信を支援する。", GOLD, GOLD_SOFT),
        ("金融・医療など", "誤情報の影響が大きい分野で、根拠確認の優先順位づけに使う。", ORANGE, ORANGE_SOFT),
    ]
    for i, (title, body, color, fill) in enumerate(uses):
        left = Inches(0.78 + (i % 2) * 6.05)
        top = Inches(1.55 + (i // 2) * 1.78)
        add_impact_card(slide, left, top, Inches(5.55), Inches(1.32), title, body, color, fill)
    add_panel(slide, Inches(1.08), Inches(5.58), Inches(11.15), Inches(0.82), RED_SOFT)
    add_text(slide, Inches(1.35), Inches(5.86), Inches(10.6), Inches(0.25), "注意: 判定結果をそのまま公開判断に使わず、人間が根拠を確認して最終判断する。", 16, RED, True, PP_ALIGN.CENTER)
    add_footer(slide, 5)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "まとめ", "Ver4の価値は、根拠付きで確認作業を速くすること")
    add_panel(slide, Inches(1.05), Inches(1.58), Inches(11.1), Inches(1.2), MINT)
    add_text(slide, Inches(1.32), Inches(1.92), Inches(10.58), Inches(0.34), "大量情報の一次スクリーニング + 根拠付きレポート化", 24, NAVY, True, PP_ALIGN.CENTER)
    add_bullets(
        slide,
        Inches(1.45),
        Inches(3.35),
        Inches(10.35),
        Inches(1.85),
        [
            "要注意情報を早く見つけ、人間の確認対象を絞り込む",
            "判定理由と確認リンクにより、調査の初動を短縮する",
            "説明責任を支え、広報・危機管理の対応を早める",
            "AIは最終判断者ではなく、人間の判断を支援する位置づけ",
        ],
        17,
    )
    add_text(slide, Inches(1.15), Inches(6.14), Inches(11.0), Inches(0.38), "プレゼンでの一言: Ver4は、人がより早く、根拠を持って確認するための支援ツールです。", 17, TEAL, True, PP_ALIGN.CENTER)
    add_footer(slide, 6)

    prs.save(OUTPUT_PPTX)
    OUTPUT_NOTES.write_text(build_notes(), encoding="utf-8")


if __name__ == "__main__":
    build_presentation()
    print(f"saved: {OUTPUT_PPTX}")
    print(f"saved: {OUTPUT_NOTES}")
