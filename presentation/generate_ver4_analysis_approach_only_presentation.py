from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PRESENTATION_DIR = Path(__file__).resolve().parent
OUTPUT_PPTX = PRESENTATION_DIR / "ver4_analysis_approach_only_20260429.pptx"
OUTPUT_NOTES = PRESENTATION_DIR / "ver4_analysis_approach_only_20260429_notes.md"

FONT = "Yu Gothic"
BG = RGBColor(247, 248, 246)
INK = RGBColor(27, 35, 43)
SUBTLE = RGBColor(88, 99, 110)
WHITE = RGBColor(255, 255, 255)
NAVY = RGBColor(24, 48, 72)
TEAL = RGBColor(15, 124, 128)
GREEN = RGBColor(61, 143, 99)
GOLD = RGBColor(191, 151, 48)
ORANGE = RGBColor(207, 115, 47)
RED = RGBColor(182, 63, 63)
LINE = RGBColor(214, 220, 224)
MINT = RGBColor(230, 243, 240)
BLUE_SOFT = RGBColor(232, 239, 248)
GREEN_SOFT = RGBColor(232, 244, 235)
GOLD_SOFT = RGBColor(251, 246, 226)
ORANGE_SOFT = RGBColor(250, 238, 227)
RED_SOFT = RGBColor(249, 232, 232)


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_text(slide, left, top, width, height, text="", size=20, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(3)
    tf.margin_right = Pt(3)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if text:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        p.alignment = align
    return box


def add_header(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, Inches(0.62), Inches(0.34), Inches(10.4), Inches(0.42), title, 25, NAVY, True)
    if subtitle:
        add_text(slide, Inches(0.64), Inches(0.81), Inches(11.1), Inches(0.26), subtitle, 10, SUBTLE)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(1.1), Inches(12.05), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def add_footer(slide, page: int) -> None:
    add_text(slide, Inches(0.65), Inches(7.05), Inches(7.2), Inches(0.18), "Fake News Checker Ver4 | 分析アプローチ", 8, SUBTLE)
    add_text(slide, Inches(12.25), Inches(7.05), Inches(0.55), Inches(0.18), str(page), 8, SUBTLE, False, PP_ALIGN.RIGHT)


def add_panel(slide, left, top, width, height, fill=WHITE, line=LINE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_bullets(slide, left, top, width, height, items: list[str], size=17, color=INK):
    box = add_text(slide, left, top, width, height)
    tf = box.text_frame
    tf.clear()
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
        p.line_spacing = 1.12


def add_step(slide, left, top, width, height, title: str, body: str, accent=TEAL, number: str | None = None):
    add_panel(slide, left, top, width, height)
    if number:
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.13), top + Inches(0.15), Inches(0.34), Inches(0.34))
        badge.fill.solid()
        badge.fill.fore_color.rgb = accent
        badge.line.fill.background()
        add_text(slide, left + Inches(0.13), top + Inches(0.2), Inches(0.34), Inches(0.16), number, 9, WHITE, True, PP_ALIGN.CENTER)
        title_left = left + Inches(0.55)
    else:
        title_left = left + Inches(0.18)
    add_text(slide, title_left, top + Inches(0.15), width - Inches(0.68), Inches(0.24), title, 13, accent, True)
    add_text(slide, left + Inches(0.18), top + Inches(0.5), width - Inches(0.35), height - Inches(0.58), body, 10, INK)


def add_arrow(slide, left, top, width, height, color=SUBTLE):
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()


def build_notes() -> str:
    return "\n".join(
        [
            "# Ver4 Analysis Approach Only Notes",
            "",
            "- generated_on: 2026-04-29",
            "",
            "## Slide 1",
            "- Ver4の分析アプローチだけを説明する抜粋資料。",
            "",
            "## Slide 2",
            "- 初期仮説は、文章だけでは真偽判定が安定しないというもの。",
            "- 外部根拠確認を入れる理由を説明する。",
            "",
            "## Slide 3",
            "- ローカル一次判定、GPT一次レビュー、Gemini外部根拠確認、補正ルールの流れを説明する。",
            "",
            "## Slide 4",
            "- 5段階ラベルの考え方を説明する。",
            "",
            "## Slide 5",
            "- 判定ロジックのキャリブレーション方針を説明する。",
            "- モデル再学習ではなく、LLM出力と根拠確認結果の解釈を改善した点を強調する。",
        ]
    ) + "\n"


def build_presentation() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    add_text(slide, Inches(0.82), Inches(1.32), Inches(10), Inches(0.42), "Fake News Checker Ver4", 25, RGBColor(205, 229, 227), True)
    add_text(slide, Inches(0.82), Inches(1.9), Inches(10.8), Inches(1.05), "分析アプローチ", 48, WHITE, True)
    add_text(slide, Inches(0.86), Inches(3.12), Inches(9.1), Inches(0.55), "文章理解 × 外部根拠確認 × 5段階判定", 20, RGBColor(232, 238, 244))
    add_panel(slide, Inches(0.88), Inches(4.58), Inches(10.4), Inches(0.82), RGBColor(37, 67, 96), RGBColor(37, 67, 96))
    add_text(slide, Inches(1.1), Inches(4.86), Inches(9.9), Inches(0.24), "文章だけで決め打ちせず、根拠と照合してから判定する", 18, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(0.88), Inches(6.42), Inches(4.2), Inches(0.22), "2026-04-29", 12, RGBColor(214, 224, 232))

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "初期仮説", "文章だけでは、ニュースの真偽判定は安定しない")
    add_text(slide, Inches(0.78), Inches(1.55), Inches(5.5), Inches(0.36), "なぜ文章だけでは難しいのか", 23, NAVY, True)
    add_bullets(
        slide,
        Inches(0.85),
        Inches(2.15),
        Inches(5.2),
        Inches(2.4),
        [
            "数字・時期・固有名詞の小さなズレが判定を左右する",
            "一部だけ正しい主張は、単純な真偽分類では扱いにくい",
            "根拠不足や学説未確立の話題は、断定すると誤判定になりやすい",
        ],
        16,
    )
    add_panel(slide, Inches(6.6), Inches(1.55), Inches(5.65), Inches(3.1), BLUE_SOFT)
    add_text(slide, Inches(6.92), Inches(1.9), Inches(5.0), Inches(0.3), "Ver4の方針", 21, TEAL, True)
    add_text(
        slide,
        Inches(6.95),
        Inches(2.48),
        Inches(4.85),
        Inches(1.05),
        "LLMの文脈理解だけに頼らず、外部根拠との照合を入れて、最終判定を5段階に整理する。",
        20,
        INK,
    )
    add_text(slide, Inches(0.92), Inches(5.62), Inches(11.4), Inches(0.4), "ポイント: 「それっぽい文章」ではなく「確認できる根拠」と照合する。", 20, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 2)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "分析パイプライン", "入力から最終ラベルまでを段階的に確認する")
    steps = [
        ("入力", "記事本文・主張・確認対象を受け取る", TEAL),
        ("ローカル一次判定", "注意度や候補ラベルを先に整理", GREEN),
        ("GPT一次レビュー", "文脈理解と論点抽出を行う", GOLD),
        ("Gemini根拠確認", "外部情報と主張を照合する", ORANGE),
        ("補正ルール", "5段階ラベルへキャリブレーション", RED),
        ("出力", "判定・理由・確認リンクを返す", NAVY),
    ]
    x0 = Inches(0.55)
    y0 = Inches(2.0)
    w = Inches(1.72)
    h = Inches(1.25)
    gap = Inches(0.28)
    for i, (title, body, color) in enumerate(steps):
        left = x0 + i * (w + gap)
        add_step(slide, left, y0, w, h, title, body, color, str(i + 1))
        if i < len(steps) - 1:
            add_arrow(slide, left + w + Inches(0.04), y0 + Inches(0.48), Inches(0.22), Inches(0.25), SUBTLE)
    add_panel(slide, Inches(1.05), Inches(4.55), Inches(11.25), Inches(1.34), MINT)
    add_text(slide, Inches(1.35), Inches(4.86), Inches(10.65), Inches(0.34), "LLMの出力をそのまま採用せず、根拠確認結果と合わせて最終判断する", 22, NAVY, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(1.45), Inches(5.38), Inches(10.35), Inches(0.26), "この補正部分が、Ver4の判定ロジックの中心。", 14, SUBTLE, False, PP_ALIGN.CENTER)
    add_footer(slide, 3)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "5段階判定", "現実の主張にあるグラデーションを表現する")
    labels = [
        ("正確", "主要部分が根拠と一致している", GREEN, GREEN_SOFT),
        ("ほぼ正確", "大筋は正しいが、数字・時期・表現などに軽微なズレがある", TEAL, MINT),
        ("判断保留", "根拠不足、未確定、または断定できない", GOLD, GOLD_SOFT),
        ("不正確", "一部に事実はあるが、重要部分に誤りがある", ORANGE, ORANGE_SOFT),
        ("誤り", "核心部分が根拠によって明確に否定される", RED, RED_SOFT),
    ]
    for i, (label, desc, color, fill) in enumerate(labels):
        top = Inches(1.55 + i * 0.86)
        add_panel(slide, Inches(0.82), top, Inches(11.78), Inches(0.62), fill)
        add_text(slide, Inches(1.1), top + Inches(0.14), Inches(1.75), Inches(0.22), label, 17, color, True)
        add_text(slide, Inches(3.02), top + Inches(0.14), Inches(8.85), Inches(0.22), desc, 16, INK)
    add_text(slide, Inches(0.92), Inches(6.24), Inches(11.35), Inches(0.4), "二値分類ではなく、中間ケースや根拠不足を明示的に扱う。", 19, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 4)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "判定ロジックのキャリブレーション", "モデル再学習ではなく、LLM出力の解釈ルールを改善")
    rules = [
        ("核心が正しい + 軽微差分", "数字・時期・表現差は、必要に応じて「ほぼ正確」へ", GREEN, GREEN_SOFT),
        ("一部真実 + 重要な誤り", "正しい要素があっても、結論を歪める場合は「不正確」へ", ORANGE, ORANGE_SOFT),
        ("存在しない発言・明確な反証", "根拠で否定できる場合は「誤り」へ", RED, RED_SOFT),
        ("論争中・根拠不足", "歴史論争や未確定情報は「判断保留」へ", GOLD, GOLD_SOFT),
    ]
    for i, (title, body, color, fill) in enumerate(rules):
        left = Inches(0.78 + (i % 2) * 6.05)
        top = Inches(1.58 + (i // 2) * 1.82)
        add_panel(slide, left, top, Inches(5.55), Inches(1.38), fill)
        add_text(slide, left + Inches(0.22), top + Inches(0.22), Inches(5.05), Inches(0.28), title, 17, color, True)
        add_text(slide, left + Inches(0.22), top + Inches(0.68), Inches(5.05), Inches(0.45), body, 13, INK)
    add_panel(slide, Inches(1.16), Inches(5.64), Inches(10.95), Inches(0.72), BLUE_SOFT)
    add_text(slide, Inches(1.35), Inches(5.87), Inches(10.55), Inches(0.24), "発表での言い方: GPT/Geminiの出力を5段階ラベルへ変換する判定ロジックをキャリブレーションした。", 15, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, 5)

    prs.save(OUTPUT_PPTX)
    OUTPUT_NOTES.write_text(build_notes(), encoding="utf-8")


if __name__ == "__main__":
    build_presentation()
    print(f"saved: {OUTPUT_PPTX}")
    print(f"saved: {OUTPUT_NOTES}")
